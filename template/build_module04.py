from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "04"
MOD_TITLE = "Tables & Data Management"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), ORANGE)
tb(s, "MODULE 04", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Tables &\nData Management", Inches(0.55), Inches(1.65),
   Inches(12.2), Inches(1.6), fs=46, bold=True, color=WHITE)
tb(s, "Design tables with the right data types, enforce data integrity with constraints,\nand write INSERT, UPDATE, DELETE, and SELECT statements.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=20, color=ORANGE)
chips = [
    ("Data Types & Design",   BLUE_ACC, Inches(0.55)),
    ("Constraints",            GREEN_OK, Inches(4.45)),
    ("CRUD & SELECT",          STEEL,   Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(5.1), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.17), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 04 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (ORANGE,   "Choose appropriate SQL Server data types for common column requirements"),
    (BLUE_ACC, "Create tables with primary keys, foreign keys, unique, and NOT NULL constraints"),
    (GREEN_OK, "Insert, update, and delete rows using T-SQL DML statements"),
    (STEEL,    "Write SELECT queries with WHERE, ORDER BY, and basic JOINs"),
    (ORANGE,   "Verify data integrity by testing constraint violations at the T-SQL prompt"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Data Types ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "Table Design & Data Types", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Choosing the right data type prevents wasted space and prevents invalid data from entering the system.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
categories = [
    (NAVY,     "INTEGERS",
     [("INT",       "Standard whole number — most common PK type"),
      ("BIGINT",    "Large integers (e.g., row counters over 2 billion)"),
      ("SMALLINT",  "2-byte integer — use when range fits"),
      ("TINYINT",   "0–255 only — status codes, flags")]),
    (BLUE_ACC, "TEXT",
     [("NVARCHAR(n)", "Unicode variable-length — default for most text"),
      ("VARCHAR(n)",  "Non-Unicode — slightly smaller storage"),
      ("NCHAR(n)",    "Fixed-length Unicode — codes, padded fields"),
      ("NVARCHAR(MAX)","Up to 2 GB — avoid in index columns")]),
    (GREEN_OK, "NUMERIC & DATE",
     [("DECIMAL(p,s)", "Exact precision — use for money, measurements"),
      ("FLOAT / REAL", "Approximate — only for scientific calculations"),
      ("DATE",         "Date only — no time component"),
      ("DATETIME2",    "High-precision datetime — preferred over DATETIME")]),
    (ORANGE,   "OTHER",
     [("BIT",         "Boolean: 0, 1, or NULL — flags and toggles"),
      ("UNIQUEIDENTIFIER", "GUID — distributed/replicated PKs"),
      ("VARBINARY(n)", "Binary data: files, hashes"),
      ("XML / JSON",   "Store structured data — query with XQuery/JSON_VALUE")]),
]
for i, (col, cat, dtypes) in enumerate(categories):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.35)
    y = Inches(1.62) + row_idx * Inches(2.35)
    rect(s, x, y, Inches(5.95), Inches(2.12), col)
    tb(s, cat, x, y + Inches(0.1), Inches(5.95), Inches(0.32),
       fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, (dtype, note) in enumerate(dtypes):
        tb(s, dtype, x + Inches(0.2), y + Inches(0.55) + j * Inches(0.38),
           Inches(1.7), Inches(0.3), fs=12, bold=True, color=WHITE)
        tb(s, note, x + Inches(1.95), y + Inches(0.55) + j * Inches(0.38),
           Inches(3.85), Inches(0.3), fs=12, color=RGBColor(0xC5, 0xD8, 0xF0))
footer(s, 3)


# ─── Slide 4: Constraints ─────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Constraints: Enforcing Data Integrity", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Constraints prevent invalid data from entering the database — they are the first line of defence.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
constraints = [
    (NAVY,     "PRIMARY KEY",
     "Uniquely identifies each row. Only one per table.\nAutomatically creates a clustered index.",
     "CONSTRAINT PK_Employee PRIMARY KEY (EmployeeID)"),
    (BLUE_ACC, "FOREIGN KEY",
     "Enforces referential integrity between tables.\nPrevents orphan rows in child tables.",
     "CONSTRAINT FK_Emp_Dept FOREIGN KEY (DeptID)\n  REFERENCES Department(DeptID)"),
    (GREEN_OK, "UNIQUE",
     "Enforces uniqueness on one or more columns.\nAllows NULLs (unlike PRIMARY KEY).",
     "CONSTRAINT UQ_Email UNIQUE (Email)"),
    (ORANGE,   "NOT NULL / CHECK",
     "NOT NULL: column must have a value.\nCHECK: column must satisfy an expression.",
     "Salary DECIMAL(10,2) NOT NULL\n  CHECK (Salary > 0)"),
]
for i, (col, title, desc, example) in enumerate(constraints):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.62) + row_idx * Inches(2.4)
    rect(s, x, y, Inches(5.95), Inches(2.18), col)
    tb(s, title, x + Inches(0.2), y + Inches(0.1), Inches(5.55), Inches(0.32),
       fs=17, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.2), y + Inches(0.5), Inches(5.55), Inches(0.52),
       fs=13, color=WHITE)
    rect(s, x + Inches(0.2), y + Inches(1.08), Inches(5.55), Inches(0.88), CODE_BG)
    tb(s, example, x + Inches(0.35), y + Inches(1.16), Inches(5.2), Inches(0.65),
       fs=11, color=RGBColor(0xA8, 0xD8, 0xFF))
footer(s, 4)


# ─── Slide 5: INSERT & UPDATE ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), NAVY)
tb(s, "INSERT & UPDATE — Adding and Changing Data", Inches(0), Inches(0.38),
   SLIDE_W, Inches(0.62), fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "INSERT", Inches(0.45), Inches(1.12), Inches(6.25), Inches(0.35),
   fs=15, bold=True, color=NAVY)
tb(s, "UPDATE", Inches(6.9), Inches(1.12), Inches(6.0), Inches(0.35),
   fs=15, bold=True, color=NAVY)

insert_sql = (
    "-- Insert a single row\n"
    "INSERT INTO Employee (FirstName, LastName, DeptID, Salary)\n"
    "VALUES ('Alice', 'Smith', 3, 75000.00);\n\n"
    "-- Insert multiple rows\n"
    "INSERT INTO Employee (FirstName, LastName, DeptID, Salary)\n"
    "VALUES ('Bob',   'Jones', 2, 62000.00),\n"
    "       ('Carol', 'Lee',   3, 80000.00);\n\n"
    "-- Insert from another table\n"
    "INSERT INTO EmployeeArchive\n"
    "SELECT * FROM Employee WHERE ResignDate < '2024-01-01';"
)
code(s, insert_sql, Inches(0.45), Inches(1.52), Inches(6.25), Inches(4.15), fs=12)

update_sql = (
    "-- Update a single column\n"
    "UPDATE Employee\n"
    "SET    Salary = Salary * 1.05\n"
    "WHERE  DeptID = 3;\n\n"
    "-- Update multiple columns\n"
    "UPDATE Employee\n"
    "SET    DeptID = 4,\n"
    "       Salary = 90000.00\n"
    "WHERE  EmployeeID = 7;\n\n"
    "-- Always verify with SELECT first!\n"
    "SELECT * FROM Employee WHERE DeptID = 3;"
)
code(s, update_sql, Inches(6.9), Inches(1.52), Inches(6.0), Inches(4.15), fs=12)

tb(s, "Tip: always run SELECT before UPDATE/DELETE to confirm target rows. Wrap in BEGIN TRAN / ROLLBACK during testing.",
   Inches(0.45), Inches(5.78), Inches(12.45), Inches(0.32), fs=13, italic=True, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: DELETE & SELECT ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "DELETE & SELECT — Removing and Querying Data", Inches(0), Inches(0.38),
   SLIDE_W, Inches(0.62), fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "DELETE without WHERE deletes ALL rows. TRUNCATE is faster but non-transactional.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True,
   color=MID_GRAY, align=PP_ALIGN.CENTER)

delete_sql = (
    "-- Delete specific rows\n"
    "DELETE FROM Employee\n"
    "WHERE  ResignDate < '2023-01-01';\n\n"
    "-- Delete all rows (keep structure)\n"
    "TRUNCATE TABLE EmployeeLog;\n\n"
    "-- Safe pattern: preview first\n"
    "BEGIN TRAN;\n"
    "  DELETE FROM Employee WHERE DeptID = 99;\n"
    "  SELECT @@ROWCOUNT;  -- check count\n"
    "ROLLBACK;"  # comment: remove ROLLBACK in production
)
code(s, delete_sql, Inches(0.45), Inches(1.55), Inches(6.25), Inches(4.0), fs=12)

select_sql = (
    "-- Basic SELECT with filter\n"
    "SELECT EmployeeID, FirstName, Salary\n"
    "FROM   Employee\n"
    "WHERE  Salary > 70000\n"
    "ORDER  BY Salary DESC;\n\n"
    "-- JOIN two tables\n"
    "SELECT e.FirstName, d.DeptName\n"
    "FROM   Employee   e\n"
    "JOIN   Department d ON e.DeptID = d.DeptID\n"
    "WHERE  d.DeptName = 'Engineering';\n\n"
    "-- Aggregate\n"
    "SELECT DeptID, AVG(Salary) AS AvgSalary\n"
    "FROM   Employee\n"
    "GROUP  BY DeptID;"
)
code(s, select_sql, Inches(6.9), Inches(1.55), Inches(6.0), Inches(4.0), fs=12)

tb(s, "DELETE", Inches(0.45), Inches(1.18), Inches(6.25), Inches(0.32),
   fs=15, bold=True, color=NAVY)
tb(s, "SELECT", Inches(6.9), Inches(1.18), Inches(6.0), Inches(0.32),
   fs=15, bold=True, color=NAVY)
tb(s, "Key clauses: WHERE (filter) · ORDER BY (sort) · JOIN (combine) · GROUP BY (aggregate) · TOP (limit rows)",
   Inches(0.45), Inches(5.65), Inches(12.45), Inches(0.32), fs=13, italic=True, color=MID_GRAY)
footer(s, 6)


# ─── Slide 7: Lab 04 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), ORANGE)
tb(s, "LAB 04", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Create Tables, Apply Constraints & Run CRUD", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=32, bold=True, color=WHITE)
tb(s, "Duration: ~60 min  ·  Tool: SSMS  ·  Uses: LabDB_TSQL from Module 03 (or recreate)",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "USE LabDB_TSQL — create the Department table (DeptID PK, DeptName NVARCHAR(100) NOT NULL)",
    "Create the Employee table with PK, FK to Department, Salary DECIMAL(10,2) CHECK > 0",
    "INSERT 3 departments and 5 employees via T-SQL VALUES",
    "Run SELECT * FROM Employee ORDER BY Salary DESC — verify rows",
    "UPDATE: give all employees in DeptID=1 a 10% salary raise",
    "Try inserting an employee with a non-existent DeptID — confirm FK error",
    "DELETE one employee — verify with SELECT COUNT(*)",
    "Run SELECT with JOIN to show employee name + department name",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.08), RGBColor(0x1A, 0x0A, 0x00))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=ORANGE)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.45) * i, Inches(0.32), Inches(0.32), ORANGE)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.45) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.45) * i, Inches(6.8), Inches(0.36),
       fs=13, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.08), RGBColor(0x1A, 0x0A, 0x00))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=ORANGE)
outcomes = [
    "Tables created with correct constraints",
    "INSERT succeeds for valid rows",
    "FK violation raises an error",
    "UPDATE changes correct rows only",
    "JOIN query returns combined result",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.70) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)
tb(s, "Clean up: tables remain for use in Module 05 security testing.",
   Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3), fs=12, italic=True, color=MID_GRAY)
footer(s, 7)


# ─── Slide 8: Module Summary ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "Module 04 Summary", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "You can now build and populate a working relational schema with correct data integrity.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
takeaways = [
    (ORANGE,   "Data types prevent bad data early",
     "Use NVARCHAR for text, DECIMAL for money, DATETIME2 for dates. Avoid FLOAT for financials."),
    (BLUE_ACC, "Constraints = first line of integrity",
     "PK ensures uniqueness. FK enforces relationships. CHECK validates ranges. NOT NULL ensures presence."),
    (GREEN_OK, "DML is always transactional",
     "Wrap risky DELETE or UPDATE in BEGIN TRAN ... ROLLBACK during development."),
    (STEEL,    "SELECT is the verification tool",
     "Run SELECT before and after every DML to confirm the right rows changed."),
    (NAVY,     "Next: Module 05",
     "Security & User Management — logins, users, server roles, and database roles."),
]
for i, (col, title, desc) in enumerate(takeaways):
    col_idx = i % 2
    row_idx = i // 2
    if i == 4:
        x = Inches(3.35)
        y = Inches(5.58)
    else:
        x = Inches(0.45) + col_idx * Inches(6.38)
        y = Inches(1.72) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(5.95), Inches(1.48), col)
    tb(s, title, x + Inches(0.2), y + Inches(0.12), Inches(5.55), Inches(0.35),
       fs=16, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.2), y + Inches(0.55), Inches(5.55), Inches(0.62),
       fs=13, color=WHITE)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-04-Tables-and-Data-Management.pptx")
prs.save(out)
print(f"Saved: {out}")
