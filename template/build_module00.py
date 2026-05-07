from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY = RGBColor(0x0D, 0x1B, 0x3E)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT = RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 9

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def tb(slide, text, left, top, width, height, fs=28, bold=False, color=WHITE, align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.word_wrap = True
    p = frame.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def footer(slide, n, total=TOTAL_SLIDES, dark=True):
    fc = MID_GRAY if not dark else RGBColor(0x4A, 0x5A, 0x7A)
    tb(slide, "Module 00: Course Introduction", Inches(0.4), SLIDE_H - Inches(0.45), Inches(5.6), Inches(0.35), fs=11, color=fc)
    tb(slide, f"{n} / {total}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.45), Inches(1.5), Inches(0.35), fs=11, color=fc, align=PP_ALIGN.RIGHT)


def format_slot(start_minutes, duration_minutes):
    end_minutes = start_minutes + duration_minutes
    start_h = start_minutes // 60
    start_m = start_minutes % 60
    end_h = end_minutes // 60
    end_m = end_minutes % 60
    return f"{start_h:02d}:{start_m:02d}-{end_h:02d}:{end_m:02d}", end_minutes


def schedule_rows(day_modules, start_minutes=9 * 60):
    rows = []
    current = start_minutes
    for module_code, module_name, duration_minutes, focus in day_modules:
        slot, current = format_slot(current, duration_minutes)
        rows.append((slot, module_code, module_name, focus, duration_minutes))
    return rows


def schedule_slide(slide_number, slide_title, subtitle, day_modules, accent):
    slide = prs.slides.add_slide(blank)
    bg(slide, WHITE)
    rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), accent)
    tb(slide, slide_title, Inches(0), Inches(0.42), SLIDE_W, Inches(0.7), fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
    tb(slide, subtitle, Inches(0), Inches(1.1), SLIDE_W, Inches(0.45), fs=16, color=MID_GRAY, align=PP_ALIGN.CENTER)

    headers = [("TIME", Inches(0.6), Inches(1.55)), ("MODULE", Inches(2.15), Inches(4.35)), ("FOCUS", Inches(6.5), Inches(3.75)), ("DURATION", Inches(10.55), Inches(1.9))]
    for header_label, header_left, header_width in headers:
        rect(slide, header_left, Inches(1.72), header_width, Inches(0.45), DARK_TEXT)
        tb(slide, header_label, header_left, Inches(1.82), header_width, Inches(0.25), fs=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

    row_y = Inches(2.22)
    row_h = Inches(0.78 if len(day_modules) <= 5 else 0.62)
    body_fs = 15 if len(day_modules) <= 5 else 13
    focus_fs = 14 if len(day_modules) <= 5 else 12
    for row_index, (slot, code, module_title, focus, duration_minutes) in enumerate(schedule_rows(day_modules)):
        fill = LIGHT_BG if row_index % 2 == 0 else WHITE
        rect(slide, Inches(0.6), row_y, Inches(11.85), row_h, fill)
        rect(slide, Inches(0.6), row_y, Inches(1.55), row_h, accent if row_index % 2 == 0 else NAVY)
        tb(slide, slot, Inches(0.6), row_y + Inches(0.13), Inches(1.55), Inches(0.25), fs=body_fs, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        tb(slide, code, Inches(2.3), row_y + Inches(0.09), Inches(1.0), Inches(0.25), fs=body_fs, bold=True, color=accent if row_index % 2 == 0 else NAVY, align=PP_ALIGN.CENTER)
        tb(slide, module_title, Inches(3.15), row_y + Inches(0.06), Inches(3.3), Inches(0.46), fs=body_fs, bold=True, color=DARK_TEXT)
        tb(slide, focus, Inches(6.5), row_y + Inches(0.06), Inches(3.75), Inches(0.45), fs=focus_fs, color=MID_GRAY)
        tb(slide, f"{duration_minutes // 60 if duration_minutes % 60 == 0 else duration_minutes / 60:g} h", Inches(10.7), row_y + Inches(0.13), Inches(1.55), Inches(0.25), fs=body_fs, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
        row_y += row_h + Inches(0.09)

    tb(slide, "Instructional blocks shown sequentially from 09:00. Breaks and lunch can be inserted by the trainer between modules.", Inches(0.7), Inches(6.62), Inches(12.0), Inches(0.35), fs=13, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
    footer(slide, slide_number, dark=False)


day1_modules = [
    ("M01", "SQL Server Overview & Architecture", 60, "Editions, components, engine basics, authentication modes"),
    ("M02", "Installing & Configuring SQL Server", 120, "Planning, setup flow, Configuration Manager, TCP/IP"),
    ("M03", "Database Creation & Management", 120, "Files, logs, SSMS and T-SQL database creation"),
    ("M04", "Tables & Data Management", 120, "Design basics, data types, constraints, CRUD"),
    ("M05", "Security & User Management", 90, "Logins, users, roles, access control"),
]

day2_modules = [
    ("M06", "Backup & Restore", 120, "Backup types, recovery models, restore workflow"),
    ("M07", "Database Maintenance", 90, "Index maintenance, statistics, DBCC CHECKDB"),
    ("M08", "SQL Server Agent & Automation", 90, "Jobs, schedules, alerts, recurring admin tasks"),
    ("M09", "Monitoring & Performance Basics", 90, "Activity Monitor, DMVs, blocking and deadlocks"),
    ("M10", "Basic Troubleshooting", 60, "Logs, login failures, connection issues, bottlenecks"),
    ("M11", "Data Import/Export", 60, "Wizard-driven CSV and Excel movement, ETL basics"),
]


s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Microsoft SQL Server\nAdministration", Inches(0.8), Inches(1.3), Inches(11.7), Inches(1.6), fs=42, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Two days to install, secure, maintain, and troubleshoot a working SQL Server environment.", Inches(1.0), Inches(3.2), Inches(11.33), Inches(0.75), fs=22, color=ORANGE, align=PP_ALIGN.CENTER)
chips = [("Instructor-led", ORANGE, Inches(1.85)), ("Hands-on labs", BLUE_ACC, Inches(4.85)), ("SQL Server 2019/2022", GREEN_OK, Inches(7.95)), ("SSMS", ORANGE, Inches(10.25))]
for chip_label, chip_color, chip_left in chips:
    rect(s, chip_left, Inches(5.0), Inches(2.2), Inches(0.68), chip_color)
    tb(s, chip_label, chip_left, Inches(5.08), Inches(2.2), Inches(0.45), fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "MODULE 00  ·  COURSE INTRODUCTION  ·  OPENING SESSION", Inches(0), Inches(6.68), SLIDE_W, Inches(0.35), fs=14, color=MID_GRAY, align=PP_ALIGN.CENTER)
footer(s, 1)


s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "Course At A Glance", Inches(0), Inches(0.45), SLIDE_W, Inches(0.7), fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Set the room, align expectations, and show the full two-day operating picture before Module 1 begins.", Inches(0.7), Inches(1.15), Inches(12.0), Inches(0.45), fs=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
cards = [
    (NAVY, "AUDIENCE", ["Junior to mid-level DBAs", "IT staff supporting SQL Server", "Admins building operational confidence"]),
    (BLUE_ACC, "ENVIRONMENT", ["SQL Server 2019/2022", "SQL Server Management Studio", "Local install or training VM"]),
    (ORANGE, "COURSE SHAPE", ["2 training days", "11 modules", "Hands-on labs in every core topic"]),
]
for card_index, (card_color, card_title, lines) in enumerate(cards):
    card_left = Inches(0.55) + card_index * Inches(4.1)
    rect(s, card_left, Inches(2.0), Inches(3.8), Inches(4.3), card_color)
    tb(s, card_title, card_left, Inches(2.2), Inches(3.8), Inches(0.35), fs=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, line in enumerate(lines):
        tb(s, f"- {line}", card_left + Inches(0.18), Inches(2.82) + Inches(0.90) * j, Inches(3.15), Inches(0.5), fs=15, color=WHITE)
footer(s, 2, dark=False)


s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), NAVY)
tb(s, "Training Agenda", Inches(0), Inches(0.42), SLIDE_W, Inches(0.75), fs=40, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Day 1 establishes the platform. Day 2 keeps that platform healthy and observable.", Inches(0), Inches(1.15), SLIDE_W, Inches(0.4), fs=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
rect(s, Inches(0.7), Inches(1.95), Inches(5.8), Inches(3.9), NAVY)
rect(s, Inches(6.85), Inches(1.95), Inches(5.8), Inches(3.9), BLUE_ACC)
tb(s, "DAY 1", Inches(0.7), Inches(2.06), Inches(5.8), Inches(0.42), fs=20, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "DAY 2", Inches(6.85), Inches(2.06), Inches(5.8), Inches(0.42), fs=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for idx, (_, item, _, _) in enumerate(day1_modules):
    tb(s, item, Inches(1.08), Inches(2.62) + Inches(0.54) * idx, Inches(5.1), Inches(0.35), fs=16, color=WHITE)
for idx, (_, item, _, _) in enumerate(day2_modules):
    tb(s, item, Inches(7.18), Inches(2.52) + Inches(0.42) * idx, Inches(5.1), Inches(0.32), fs=15, color=WHITE)
tb(s, "11 technical modules across installation, security, maintenance, automation, monitoring, troubleshooting, and data movement.", Inches(0), Inches(6.35), SLIDE_W, Inches(0.4), fs=15, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
footer(s, 3, dark=False)


schedule_slide(4, "Day 1 Timetable", "Installation, configuration, database fundamentals, and security", day1_modules, ORANGE)
schedule_slide(5, "Day 2 Timetable", "Backup, maintenance, automation, monitoring, troubleshooting, and data movement", day2_modules, BLUE_ACC)


s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "Training Outcomes", Inches(0), Inches(0.45), SLIDE_W, Inches(0.7), fs=40, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
outcomes = [
    (NAVY, "INSTALL", ["Plan a SQL Server setup", "Configure services and protocols", "Connect locally and remotely"]),
    (BLUE_ACC, "BUILD", ["Create databases", "Work with tables and data", "Use SSMS and T-SQL"]),
    (ORANGE, "SECURE", ["Create logins and users", "Assign server and database roles", "Verify access control"]),
    (GREEN_OK, "OPERATE", ["Back up and restore", "Run integrity and maintenance checks", "Automate recurring tasks"]),
    (NAVY, "OBSERVE", ["Monitor activity", "Read logs", "Troubleshoot basic failures"]),
]
positions = [Inches(0.45), Inches(2.95), Inches(5.45), Inches(7.95), Inches(10.45)]
for (outcome_color, outcome_title, lines), outcome_left in zip(outcomes, positions):
    rect(s, outcome_left, Inches(1.85), Inches(2.4), Inches(4.85), outcome_color)
    tb(s, outcome_title, outcome_left, Inches(2.05), Inches(2.4), Inches(0.35), fs=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, line in enumerate(lines):
        tb(s, f"- {line}", outcome_left + Inches(0.12), Inches(2.63) + Inches(0.90) * j, Inches(2.1), Inches(0.50), fs=14, color=WHITE)
footer(s, 6, dark=False)


s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "How We Will Work", Inches(0), Inches(0.45), SLIDE_W, Inches(0.7), fs=40, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
rules = [
    "Instructor-led explanation first. Hands-on practice immediately after.",
    "Use SSMS for admin workflows and T-SQL where the module calls for it.",
    "Each lab ends with a visible verification check: connection, object, backup file, job history, or result set.",
    "Work on a safe training instance or VM. Repeatable tasks beat risky shortcuts.",
    "Ask about operational tradeoffs early: authentication mode, recovery model, permissions, maintenance timing.",
]
for i, rule in enumerate(rules):
    y = Inches(1.72) + Inches(0.90) * i
    rect(s, Inches(1.0), y, Inches(0.42), Inches(0.42), ORANGE if i % 2 == 0 else BLUE_ACC)
    tb(s, rule, Inches(1.62), y - Inches(0.02), Inches(10.3), Inches(0.65), fs=18, color=DARK_TEXT)
tb(s, "Verification-first labs keep the class aligned even when the environment differs slightly between machines.", Inches(0), Inches(6.35), SLIDE_W, Inches(0.35), fs=15, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
footer(s, 7, dark=False)


s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "What Success Looks Like", Inches(0), Inches(0.48), SLIDE_W, Inches(0.55), fs=30, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "By the end of the course, each participant should be able to operate a clean SQL Server admin workflow from setup through recovery and troubleshooting.", Inches(1.0), Inches(1.35), Inches(11.3), Inches(0.85), fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
steps = [
    (Inches(0.95), "SET UP", ["Install", "Configure", "Connect"]),
    (Inches(3.45), "BUILD", ["Create DBs", "Manage tables", "Secure access"]),
    (Inches(5.95), "PROTECT", ["Backup", "Restore", "Check integrity"]),
    (Inches(8.45), "AUTOMATE", ["Jobs", "Schedules", "Maintenance"]),
    (Inches(10.95), "OBSERVE", ["Monitor", "Investigate", "Import/Export"]),
]
for step_left, step_title, items in steps:
    rect(s, step_left, Inches(3.55), Inches(1.8), Inches(2.05), BLUE_ACC if step_title in {"BUILD", "AUTOMATE"} else ORANGE)
    tb(s, step_title, step_left, Inches(3.75), Inches(1.8), Inches(0.28), fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, item in enumerate(items):
        tb(s, item, step_left, Inches(4.18) + Inches(0.38) * j, Inches(1.8), Inches(0.25), fs=15, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 8)


s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.12), ORANGE)
tb(s, "Start with the instance and the engine.", Inches(0), Inches(1.72), SLIDE_W, Inches(0.9), fs=46, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Module 01 begins with the language and architecture every administrator needs before installation and operational work make sense.", Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.8), fs=22, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "Next: SQL Server Overview & Architecture", Inches(0), Inches(5.15), SLIDE_W, Inches(0.55), fs=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Open SSMS and be ready to identify the instance, the engine, and the system databases.", Inches(0), Inches(6.0), SLIDE_W, Inches(0.35), fs=16, color=MID_GRAY, align=PP_ALIGN.CENTER)
footer(s, 9)


repo_root = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-00-Course-Introduction.pptx")
prs.save(out)
print(f"Saved: {out}")