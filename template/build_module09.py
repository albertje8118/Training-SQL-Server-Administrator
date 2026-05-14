from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
RED_ERR  = RGBColor(0xE0, 0x52, 0x52)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "09"
MOD_TITLE = "Monitoring & Performance Basics"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), STEEL)
tb(s, "MODULE 09", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Monitoring &\nPerformance Basics", Inches(0.55), Inches(1.65),
   Inches(12.2), Inches(1.6), fs=46, bold=True, color=WHITE)
tb(s, "Find out what SQL Server is doing right now — and what is slowing it down —\n"
      "using Activity Monitor, DMVs, and execution plans.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=19, color=BLUE_ACC)
chips = [
    ("Activity Monitor",  ORANGE,   Inches(0.55)),
    ("DMVs",              GREEN_OK, Inches(4.45)),
    ("Blocking & Waits",  AMBER,    Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(4.95), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.02), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Day 2  ·  Afternoon  ·  1.5 hours", Inches(0.55), Inches(6.72), Inches(6.0), Inches(0.35),
   fs=13, color=MID_GRAY)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), STEEL)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 09 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (STEEL,    "Use Activity Monitor to identify active sessions and expensive queries"),
    (BLUE_ACC, "Query key DMVs to see CPU, memory, I/O, and wait statistics"),
    (ORANGE,   "Identify a blocking chain and find the head blocker"),
    (GREEN_OK, "Explain what wait statistics reveal about performance bottlenecks"),
    (STEEL,    "Read a basic query execution plan to find index scan warnings"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Activity Monitor ───────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Activity Monitor", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Activity Monitor gives a real-time dashboard of instance health — no T-SQL required.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

panels = [
    (NAVY,     "OVERVIEW",
     "4 live charts: % CPU, Waiting Tasks, Database I/O (MB/s), Batch Requests/sec\n"
     "Refreshes every 10 seconds by default (configurable)."),
    (BLUE_ACC, "PROCESSES",
     "Lists all active SPID sessions.\n"
     "Columns: Session ID, CPU, Memory, Blocked By, Command, Wait Type, Host.\n"
     "Right-click a session → Kill Process (use carefully!)"),
    (GREEN_OK, "RESOURCE WAITS",
     "Cumulative wait categories since last reset.\n"
     "Helps identify whether bottleneck is CPU, I/O, locks, network, or memory.\n"
     "Higher wait ms = bigger contributor."),
    (ORANGE,   "DATA FILE I/O",
     "Per-database file read/write throughput.\n"
     "Useful for identifying which database is driving I/O.\n"
     "Check against storage capacity baseline."),
    (STEEL,    "RECENT EXPENSIVE QUERIES",
     "Top queries by CPU, elapsed time, logical reads, or physical reads.\n"
     "Click a row to see the query text and execution plan.\n"
     "First place to look when 'the server is slow'."),
]
for i, (col, title, desc) in enumerate(panels):
    col_idx = i % 2
    row_idx = i // 2
    if i == 4:
        x = Inches(3.35)
        y = Inches(5.32)
    else:
        x = Inches(0.45) + col_idx * Inches(6.38)
        y = Inches(1.62) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(0.35), Inches(1.45), col)
    tb(s, title, x + Inches(0.52), y + Inches(0.08), Inches(5.55), Inches(0.32),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.52), y + Inches(0.48), Inches(5.55), Inches(0.72),
       fs=13, color=MID_GRAY)

tb(s, "Open:  SSMS toolbar → Activity Monitor icon  — or — right-click instance → Activity Monitor",
   Inches(0.45), Inches(6.72), Inches(12.45), Inches(0.28), fs=12, italic=True, color=DARK_TEXT)
footer(s, 3)


# ─── Slide 4: Key DMVs ───────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Essential DMVs for Every DBA", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Dynamic Management Views surface real-time engine data — far more detail than Activity Monitor.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

dmv_sql = (
    "-- Active sessions with wait info\n"
    "SELECT r.session_id, r.status, r.cpu_time, r.total_elapsed_time,\n"
    "       r.wait_type, r.wait_time, t.text AS sql_text\n"
    "FROM sys.dm_exec_requests r\n"
    "CROSS APPLY sys.dm_exec_sql_text(r.sql_handle) t\n"
    "WHERE r.session_id > 50  -- exclude system sessions\n"
    "ORDER BY r.cpu_time DESC;\n\n"
    "-- Top wait types since last restart\n"
    "SELECT TOP 10 wait_type, wait_time_ms, waiting_tasks_count\n"
    "FROM sys.dm_os_wait_stats\n"
    "WHERE wait_type NOT IN ('SLEEP_TASK','BROKER_TO_FLUSH','WAITFOR',...)\n"
    "ORDER BY wait_time_ms DESC;\n\n"
    "-- Memory usage by database\n"
    "SELECT database_id, DB_NAME(database_id) AS DBName,\n"
    "       COUNT(*) * 8 / 1024 AS CacheMB\n"
    "FROM sys.dm_os_buffer_descriptors\n"
    "GROUP BY database_id ORDER BY CacheMB DESC;"
)
code(s, dmv_sql, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.55), fs=11)

dmv_table = [
    ("sys.dm_exec_requests",         "Currently executing requests — CPU, waits, query text"),
    ("sys.dm_exec_sessions",         "All open sessions — login, host, program name"),
    ("sys.dm_exec_sql_text",         "Table-valued function — returns SQL text from sql_handle"),
    ("sys.dm_os_wait_stats",         "Cumulative wait type counts and durations"),
    ("sys.dm_os_buffer_descriptors", "Buffer pool pages — which DBs consume most memory"),
    ("sys.dm_io_virtual_file_stats", "Actual read/write I/O per database file"),
    ("sys.dm_exec_query_stats",      "Cached query plan stats — logical reads, CPU, executions"),
]
for i, (dmv, desc) in enumerate(dmv_table):
    col = [STEEL, BLUE_ACC, GREEN_OK, ORANGE, NAVY, STEEL, BLUE_ACC][i]
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.38), y + Inches(0.1), Inches(0.28), Inches(0.28), col)
    tb(s, dmv, Inches(7.78), y, Inches(5.18), Inches(0.28), fs=12, bold=True, color=DARK_TEXT)
    tb(s, desc, Inches(7.78), y + Inches(0.28), Inches(5.18), Inches(0.3), fs=11, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: Blocking & Deadlocks ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), AMBER)
tb(s, "Blocking & Deadlocks", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Blocking degrades throughput. Deadlocks crash sessions. Both have T-SQL diagnostic queries.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Blocking side
rect(s, Inches(0.45), Inches(1.55), Inches(5.95), Inches(2.35), NAVY)
tb(s, "BLOCKING", Inches(0.45), Inches(1.68), Inches(5.95), Inches(0.35),
   fs=20, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
block_pts = [
    "Session A holds a lock; Session B waits",
    "Blocking is normal — only a problem if > 5 sec",
    "Head blocker: the session that no-one is waiting for",
    "sys.dm_exec_requests.blocking_session_id = head blocker",
    "Resolution: commit/rollback the long transaction or kill SPID",
]
for j, pt in enumerate(block_pts):
    tb(s, f"·  {pt}", Inches(0.65), Inches(2.18) + j * Inches(0.36),
       Inches(5.55), Inches(0.28), fs=13, color=WHITE)

# Deadlock side
rect(s, Inches(6.78), Inches(1.55), Inches(6.1), Inches(2.35), RGBColor(0x4A, 0x10, 0x10))
tb(s, "DEADLOCK", Inches(6.78), Inches(1.68), Inches(6.1), Inches(0.35),
   fs=20, bold=True, color=RED_ERR, align=PP_ALIGN.CENTER)
dead_pts = [
    "Session A waits for B; Session B waits for A",
    "SQL Server detects and kills the 'deadlock victim'",
    "Victim receives error 1205 — must retry the transaction",
    "Trace flag 1222 or Extended Events to capture deadlock graph",
    "Resolution: consistent lock-acquisition order; shorter transactions",
]
for j, pt in enumerate(dead_pts):
    tb(s, f"·  {pt}", Inches(6.98), Inches(2.18) + j * Inches(0.36),
       Inches(5.7), Inches(0.28), fs=13, color=WHITE)

# Diagnostic query
tb(s, "FIND THE HEAD BLOCKER", Inches(0.45), Inches(3.98), Inches(12.45), Inches(0.32),
   fs=15, bold=True, color=DARK_TEXT)
blocking_sql = (
    "SELECT\n"
    "  r.session_id,\n"
    "  r.blocking_session_id,\n"
    "  r.wait_type,\n"
    "  r.wait_time / 1000.0 AS wait_sec,\n"
    "  t.text AS blocking_sql\n"
    "FROM sys.dm_exec_requests r\n"
    "CROSS APPLY sys.dm_exec_sql_text(r.sql_handle) t\n"
    "WHERE r.blocking_session_id > 0;"
)
code(s, blocking_sql, Inches(0.45), Inches(4.35), Inches(7.55), Inches(2.75), fs=12)

tb(s, "Kill head blocker (last resort):\n\nKILL <session_id>;",
   Inches(8.22), Inches(4.35), Inches(4.66), Inches(0.92), fs=14, color=DARK_TEXT)
code(s, "KILL 58;  -- kills session 58", Inches(8.22), Inches(5.28), Inches(4.66), Inches(0.5), fs=12)
tb(s, "Warning: KILL causes an immediate rollback — can be slow for large transactions.",
   Inches(8.22), Inches(5.88), Inches(4.66), Inches(0.65), fs=12, italic=True, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: Wait Statistics & Execution Plans ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "Wait Statistics & Execution Plans", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Wait types tell you WHY queries are slow. Execution plans show you WHERE time is spent.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

wait_types = [
    (ORANGE,   "PAGEIOLATCH_SH/_EX",  "Waiting for a data page to load from disk → I/O bottleneck"),
    (RED_ERR,  "LCK_M_X / LCK_M_S",  "Lock waits → blocking by another transaction"),
    (BLUE_ACC, "CXPACKET / CXCONSUMER","Parallel query coordination → consider MAXDOP tuning"),
    (GREEN_OK, "ASYNC_NETWORK_IO",    "Client not consuming results fast enough → app layer issue"),
    (STEEL,    "SOS_SCHEDULER_YIELD", "CPU pressure → queries spinning on scheduler"),
    (AMBER,    "WRITELOG",            "Log flush waits → storage I/O latency on log drive"),
]
rect(s, Inches(0.45), Inches(1.55), Inches(6.55), Inches(4.28), NAVY)
tb(s, "COMMON WAIT TYPES", Inches(0.45), Inches(1.68), Inches(6.55), Inches(0.32),
   fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for j, (col, wt, meaning) in enumerate(wait_types):
    y = Inches(2.08) + j * Inches(0.58)
    rect(s, Inches(0.62), y + Inches(0.05), Inches(0.22), Inches(0.28), col)
    tb(s, wt, Inches(1.02), y, Inches(2.55), Inches(0.28), fs=12, bold=True, color=WHITE)
    tb(s, meaning, Inches(3.62), y, Inches(3.2), Inches(0.3), fs=11, color=MID_GRAY)

rect(s, Inches(7.25), Inches(1.55), Inches(5.63), Inches(4.28), RGBColor(0x05, 0x12, 0x2A))
tb(s, "READING EXECUTION PLANS", Inches(7.25), Inches(1.68), Inches(5.63), Inches(0.32),
   fs=15, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
plan_items = [
    ("Show actual plan:",    "Query → Include Actual Execution Plan (Ctrl+M)"),
    ("Read right to left:",  "Data flows from rightmost operator to leftmost (SELECT)"),
    ("Thick arrows = cost:", "Wider line = more rows estimated — check for skew"),
    ("Table Scan warning:",  "Orange exclamation = no suitable index — consider adding one"),
    ("Key Lookup:",          "Goes to clustered index for columns not in non-clustered → add INCLUDE"),
    ("Cost %:",              "Percentage of total plan cost — highest % = tune first"),
]
for j, (label, detail) in enumerate(plan_items):
    y = Inches(2.08) + j * Inches(0.58)
    tb(s, label, Inches(7.42), y, Inches(1.88), Inches(0.28), fs=12, bold=True, color=BLUE_ACC)
    tb(s, detail, Inches(9.35), y, Inches(3.35), Inches(0.3), fs=11, color=WHITE)

tb(s, "Estimated vs Actual rows divergence > 10× = statistics are stale → run UPDATE STATISTICS.",
   Inches(0.45), Inches(5.95), Inches(12.45), Inches(0.3), fs=13, italic=True, color=DARK_TEXT)
footer(s, 6)


# ─── Slide 7: Lab 09 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), STEEL)
tb(s, "LAB 09", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Investigate Dian's Slow Inventory Report", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~45 min  ·  Tool: SSMS  ·  Database: NovaMart",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Open Activity Monitor — observe CPU %, Waiting Tasks, Batch Requests/sec",
    "Expand 'Recent Expensive Queries' — find the inventory report query",
    "Run the DMV query for active sessions — review wait_type column",
    "Run sys.dm_os_wait_stats — identify top 5 wait types",
    "Open a new query: SELECT * FROM Product ORDER BY ProductName (no index on name)",
    "Enable 'Include Actual Execution Plan' — run the query",
    "Read the plan: identify Table Scan warning",
    "Simulate blocking: open 2 windows, BEGIN TRAN + UPDATE in window 1, SELECT in window 2",
    "Identify the blocking session using dm_exec_requests",
    "Roll back window 1 transaction — confirm window 2 unblocks",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.45), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=STEEL)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.05) + Inches(0.41) * i, Inches(0.38), Inches(0.28), STEEL)
    tb(s, str(i + 1), Inches(0.7), Inches(3.07) + Inches(0.41) * i, Inches(0.38), Inches(0.24),
       fs=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.22), Inches(3.04) + Inches(0.41) * i, Inches(6.72), Inches(0.34),
       fs=12, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.45), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "Activity Monitor loads successfully",
    "DMV returns active session rows",
    "Top wait types identified",
    "Table Scan visible in plan",
    "Blocking chain identified via DMV",
    "Blocking resolves after ROLLBACK",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.05) + Inches(0.64) * i,
       Inches(3.88), Inches(0.48), fs=12, color=WHITE)
footer(s, 7)


# ─── Slide 8: Performance Checklist ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "DBA Performance First-Response Checklist", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "When someone says 'SQL Server is slow', start here.",
   Inches(0), Inches(1.08), SLIDE_W, Inches(0.38), fs=20, italic=True, color=BLUE_ACC,
   align=PP_ALIGN.CENTER)

checks = [
    (ORANGE,   "Check CPU & Batch Requests",     "Activity Monitor overview charts — is CPU spiking or flat?"),
    (BLUE_ACC, "Check Waiting Tasks",             "Activity Monitor — large Waiting Tasks = blocking or resource pressure"),
    (GREEN_OK, "Find Blocking Sessions",          "dm_exec_requests WHERE blocking_session_id > 0"),
    (STEEL,    "Check Top Wait Types",            "dm_os_wait_stats — top wait category guides next steps"),
    (AMBER,    "Find Expensive Queries",          "Activity Monitor Recent Expensive Queries OR dm_exec_query_stats"),
    (NAVY,     "Check Index Fragmentation",       "dm_db_index_physical_stats — run if I/O is the wait type"),
    (ORANGE,   "Check Available Memory",          "dm_os_sys_info.physical_memory_in_use_kb vs total"),
    (BLUE_ACC, "Review SQL Server Error Log",     "Management → SQL Server Logs — any recent errors or warnings?"),
]
for i, (col, title, detail) in enumerate(checks):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.62) + row_idx * Inches(1.28)
    rect(s, x, y, Inches(0.38), Inches(1.05), col)
    tb(s, title, x + Inches(0.55), y + Inches(0.08), Inches(5.55), Inches(0.32),
       fs=15, bold=True, color=WHITE)
    tb(s, detail, x + Inches(0.55), y + Inches(0.48), Inches(5.55), Inches(0.42),
       fs=12, color=MID_GRAY)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-09-Monitoring-and-Performance-Basics.pptx")
prs.save(out)
print(f"Saved: {out}")
