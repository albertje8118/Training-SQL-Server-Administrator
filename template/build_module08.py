from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "08"
MOD_TITLE = "SQL Server Agent & Automation"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 7

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), BLUE_ACC)
tb(s, "MODULE 08", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "SQL Server Agent\n& Automation", Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.6),
   fs=46, bold=True, color=WHITE)
tb(s, "Stop doing maintenance by hand — let SQL Server Agent\n"
      "run backups, index jobs, and alerts while you sleep.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=20, color=BLUE_ACC)
chips = [
    ("SQL Agent Overview", ORANGE,   Inches(0.55)),
    ("Jobs & Steps",        GREEN_OK, Inches(4.45)),
    ("Schedules & Alerts",  STEEL,    Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(4.95), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.02), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Day 2  ·  Morning  ·  1.5 hours", Inches(0.55), Inches(6.72), Inches(6.0), Inches(0.35),
   fs=13, color=MID_GRAY)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_ACC)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 08 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (BLUE_ACC, "Describe what SQL Server Agent is and list its components"),
    (ORANGE,   "Create a multi-step Agent job using the SSMS GUI"),
    (GREEN_OK, "Attach a recurring schedule to a job (nightly, weekly)"),
    (STEEL,    "Configure an operator and alert to receive failure notifications"),
    (BLUE_ACC, "View job history and diagnose a failed job step"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: SQL Server Agent Architecture ───────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "SQL Server Agent Architecture", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "SQL Server Agent is a Windows service that runs scheduled T-SQL, SSIS packages, and operating system commands.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=16, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

components = [
    (NAVY,     "JOB",
     "A named automation task.\n"
     "Contains one or more steps.\n"
     "Can be run on demand or by schedule.\n"
     "Stored in msdb..sysjobs."),
    (BLUE_ACC, "JOB STEP",
     "Individual unit of work within a job.\n"
     "Types: T-SQL, PowerShell, SSIS, OS command.\n"
     "On failure: stop, go to next step, or go to a specific step.\n"
     "Has its own success/failure action."),
    (GREEN_OK, "SCHEDULE",
     "Defines WHEN a job runs.\n"
     "Recurring: daily, weekly, every N minutes.\n"
     "One schedule can drive multiple jobs.\n"
     "Stored in msdb..sysschedules."),
    (STEEL,    "ALERT",
     "Triggered by SQL Server error number or\n"
     "performance condition threshold.\n"
     "Can execute a job and/or notify operators.\n"
     "Useful for critical error numbers (823, 824, 825)."),
    (ORANGE,   "OPERATOR",
     "Recipient of alert notifications.\n"
     "Notification via email (Database Mail).\n"
     "Requires SQL Server Agent Mail profile.\n"
     "Multiple operators can receive the same alert."),
]
for i, (col, title, desc) in enumerate(components):
    col_idx = i % 3
    row_idx = i // 3
    x = Inches(0.35) + col_idx * Inches(4.28)
    y = Inches(1.62) + row_idx * Inches(2.45)
    if i == 3:
        x = Inches(1.5)
    if i == 4:
        x = Inches(5.78)
    rect(s, x, y, Inches(3.85), Inches(2.15), col)
    tb(s, title, x + Inches(0.18), y + Inches(0.1), Inches(3.48), Inches(0.35),
       fs=16, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.18), y + Inches(0.52), Inches(3.48), Inches(1.45),
       fs=12, color=WHITE)

tb(s, "Prerequisites: SQL Server Agent service must be Running. Check in SQL Server Configuration Manager.",
   Inches(0.35), Inches(6.72), Inches(12.6), Inches(0.28), fs=12, italic=True, color=MID_GRAY)
footer(s, 3)


# ─── Slide 4: Creating a Job in SSMS ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Creating a Job in SSMS", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Walk-through: create the NovaMart Nightly Backup Job.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

steps_gui = [
    (NAVY,     "1. Open SQL Server Agent",
     "In SSMS Object Explorer → SQL Server Agent → right-click Jobs → New Job"),
    (BLUE_ACC, "2. General Tab",
     "Name: 'NovaMart Nightly Backup'  ·  Owner: sa (or a dedicated service account)  ·  Enabled: checked"),
    (GREEN_OK, "3. Steps Tab",
     "Click New → Step Name: 'Full Backup'  ·  Type: T-SQL  ·  Database: master\n"
     "Command:  BACKUP DATABASE NovaMart TO DISK='C:\\SQLBackup\\NovaMart_Full.bak' WITH FORMAT,COMPRESSION;"),
    (ORANGE,   "4. Steps — Add Log Backup Step",
     "Click New → 'Log Backup'  ·  T-SQL  ·  Database: master\n"
     "Command:  BACKUP LOG NovaMart TO DISK='C:\\SQLBackup\\NovaMart_Log.bak' WITH INIT;"),
    (STEEL,    "5. Schedules Tab",
     "Click New → Name: 'Nightly 02:00'  ·  Recurring: Daily  ·  Time: 02:00:00  ·  OK"),
    (AMBER,    "6. Notifications Tab",
     "On job failure: email an operator (if Database Mail is configured)"),
]
for i, (col, title, desc) in enumerate(steps_gui):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.35)
    y = Inches(1.62) + row_idx * Inches(1.65)
    rect(s, x, y, Inches(0.38), Inches(1.42), col)
    tb(s, title, x + Inches(0.55), y + Inches(0.08), Inches(5.55), Inches(0.32),
       fs=14, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.55), y + Inches(0.48), Inches(5.55), Inches(0.78),
       fs=12, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: Schedules, Alerts, and Monitoring Jobs ─────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "Schedules, Alerts & Job History", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Once jobs run unattended, you need reliable monitoring to catch failures early.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Schedule types
rect(s, Inches(0.45), Inches(1.55), Inches(5.85), Inches(2.35), NAVY)
tb(s, "SCHEDULE FREQUENCIES", Inches(0.45), Inches(1.68), Inches(5.85), Inches(0.35),
   fs=16, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
freqs = ["One time only", "Daily", "Weekly (day of week)", "Monthly (day / week of month)",
         "When SQL Server Agent starts", "When CPU is idle"]
for j, f in enumerate(freqs):
    tb(s, f"·  {f}", Inches(0.65), Inches(2.12) + j * Inches(0.36),
       Inches(5.45), Inches(0.28), fs=13, color=WHITE)

# Job history
rect(s, Inches(6.78), Inches(1.55), Inches(6.1), Inches(2.35), RGBColor(0x05, 0x12, 0x2A))
tb(s, "VIEWING JOB HISTORY", Inches(6.78), Inches(1.68), Inches(6.1), Inches(0.35),
   fs=16, bold=True, color=GREEN_OK, align=PP_ALIGN.CENTER)
history_items = [
    "Right-click Job → View History",
    "Green check = success  ·  Red X = failure",
    "Expand failed step to read error message",
    "msdb..sysjobhistory — query history via T-SQL",
    "Filter by date range in the Log File Viewer",
]
for j, item in enumerate(history_items):
    tb(s, f"·  {item}", Inches(6.98), Inches(2.12) + j * Inches(0.36),
       Inches(5.7), Inches(0.28), fs=13, color=WHITE)

# T-SQL job execution
tb(s, "RUN & QUERY JOBS VIA T-SQL", Inches(0.45), Inches(3.98), Inches(12.45), Inches(0.32),
   fs=16, bold=True, color=DARK_TEXT)
job_sql = (
    "-- Run a job immediately\n"
    "EXEC msdb.dbo.sp_start_job N'NovaMart Nightly Backup';\n\n"
    "-- Check job status\n"
    "SELECT j.name, ja.run_requested_date, ja.stop_execution_date,\n"
    "       jh.run_status  -- 0=failed 1=success 2=retry 3=cancelled\n"
    "FROM msdb.dbo.sysjobs j\n"
    "JOIN msdb.dbo.sysjobactivity ja ON j.job_id = ja.job_id\n"
    "JOIN msdb.dbo.sysjobhistory jh ON j.job_id = jh.job_id\n"
    "WHERE j.name = 'NovaMart Nightly Backup';"
)
code(s, job_sql, Inches(0.45), Inches(4.35), Inches(12.45), Inches(2.75), fs=12)
footer(s, 5)


# ─── Slide 6: T-SQL to Create a Job ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Creating a Job via T-SQL", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Scripting jobs allows them to be version-controlled, audited, and deployed to multiple instances.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

create_job_sql = (
    "USE msdb;\n"
    "GO\n\n"
    "-- Create the job\n"
    "EXEC sp_add_job\n"
    "  @job_name = N'NovaMart Nightly Backup';\n\n"
    "-- Add a T-SQL step\n"
    "EXEC sp_add_jobstep\n"
    "  @job_name = N'NovaMart Nightly Backup',\n"
    "  @step_name = N'Full Backup',\n"
    "  @subsystem = N'TSQL',\n"
    "  @command = N'BACKUP DATABASE NovaMart\n"
    "    TO DISK = ''C:\\SQLBackup\\NovaMart_Full.bak''\n"
    "    WITH FORMAT, COMPRESSION, STATS = 10;',\n"
    "  @database_name = N'master';\n\n"
    "-- Add a daily schedule at 02:00\n"
    "EXEC sp_add_schedule\n"
    "  @schedule_name = N'Nightly 02:00',\n"
    "  @freq_type = 4,        -- daily\n"
    "  @freq_interval = 1,\n"
    "  @active_start_time = 20000;  -- 02:00:00\n\n"
    "-- Attach schedule to job\n"
    "EXEC sp_attach_schedule\n"
    "  @job_name = N'NovaMart Nightly Backup',\n"
    "  @schedule_name = N'Nightly 02:00';\n\n"
    "-- Register the job on this server\n"
    "EXEC sp_add_jobserver\n"
    "  @job_name = N'NovaMart Nightly Backup';"
)
code(s, create_job_sql, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.55), fs=10)

annotations = [
    (ORANGE,   "sp_add_job",         "Creates the job container in msdb — just name and owner."),
    (BLUE_ACC, "sp_add_jobstep",     "Adds a step; @subsystem = TSQL, SSIS, PowerShell, CmdExec."),
    (GREEN_OK, "sp_add_schedule",    "@freq_type: 1=once, 4=daily, 8=weekly, 16=monthly."),
    (STEEL,    "@active_start_time", "Format HHMMSS as integer: 020000 = 02:00:00."),
    (AMBER,    "sp_attach_schedule", "One schedule can be shared across multiple jobs."),
    (NAVY,     "sp_add_jobserver",   "Required to register job on the local server (@server_name=N'(local)')."),
    (ORANGE,   "Script existing job", "In SSMS: right-click Job → Script Job as → CREATE To → clipboard."),
]
for i, (col, kw, note) in enumerate(annotations):
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.52), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(5.0), Inches(0.25),
       fs=13, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.30), Inches(5.0), Inches(0.28),
       fs=12, color=MID_GRAY)
footer(s, 6)


# ─── Slide 7: Lab 08 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), BLUE_ACC)
tb(s, "LAB 08", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Automate NovaMart Nightly Maintenance", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~45 min  ·  Tool: SSMS  ·  Database: msdb + NovaMart",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Open SSMS → confirm SQL Server Agent service is Running",
    "Create a new job: 'NovaMart Nightly Backup' (owner: sa)",
    "Step 1: 'Full Backup' — T-SQL — BACKUP DATABASE NovaMart ... WITH FORMAT, COMPRESSION",
    "Step 2: 'Update Stats' — T-SQL — EXEC sp_updatestats",
    "Step 3: 'Log Backup' — T-SQL — BACKUP LOG NovaMart ... WITH INIT",
    "Add schedule: Daily at 02:00 — name: 'Nightly 02:00'",
    "Right-click job → Start Job at Step → Step 1 — run manually to test",
    "View job history — confirm all 3 steps show green checkmark",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=BLUE_ACC)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.47) * i, Inches(0.32), Inches(0.32), BLUE_ACC)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.47) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.47) * i, Inches(6.8), Inches(0.36),
       fs=12, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "Job visible under SQL Agent → Jobs",
    "All 3 steps complete successfully",
    "Backup file created on disk",
    "Job history shows duration & status",
    "Schedule appears in Job Properties",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.70) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)
footer(s, 7)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-08-SQL-Server-Agent-and-Automation.pptx")
prs.save(out)
print(f"Saved: {out}")
