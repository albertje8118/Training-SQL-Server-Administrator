from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "11"
MOD_TITLE = "Data Import & Export"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), GREEN_OK)
tb(s, "MODULE 11", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Data Import & Export", Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.2),
   fs=46, bold=True, color=WHITE)
tb(s, "Move data in and out of SQL Server using the Import/Export Wizard,\n"
      "BULK INSERT, BCP, and basic ETL concepts.",
   Inches(0.55), Inches(3.35), Inches(11.0), Inches(0.9), fs=20, color=GREEN_OK)
chips = [
    ("Import/Export Wizard", BLUE_ACC, Inches(0.55)),
    ("BULK INSERT / BCP",     ORANGE,   Inches(4.45)),
    ("ETL Concepts",          STEEL,    Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(4.95), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.02), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Day 2  ·  Afternoon  ·  1 hour", Inches(0.55), Inches(6.72), Inches(6.0), Inches(0.35),
   fs=13, color=MID_GRAY)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN_OK)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 11 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (GREEN_OK, "Use the SSMS Import/Export Wizard to import a CSV file into a table"),
    (BLUE_ACC, "Import data from a flat file using BULK INSERT with format options"),
    (ORANGE,   "Export a SQL Server table to Excel using the Export Wizard"),
    (STEEL,    "Use the BCP utility from the command line to import and export data"),
    (GREEN_OK, "Explain the Extract, Transform, Load (ETL) pattern and where SSIS fits"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Import/Export Wizard ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "SQL Server Import/Export Wizard", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "A GUI-driven SSIS package generator — easiest path for one-off imports and exports.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

steps_wizard = [
    (NAVY,     "1. Launch the Wizard",
     "SSMS → right-click database → Tasks → Import Data  (or Export Data)\n"
     "Or: Start menu → SQL Server Import and Export Wizard"),
    (BLUE_ACC, "2. Choose Data Source",
     "For CSV import: select 'Flat File Source'\n"
     "Browse to file → preview → set delimiter (comma, tab, semicolon)\n"
     "Check 'Column names in first data row' if applicable"),
    (GREEN_OK, "3. Choose Destination",
     "Select 'SQL Server Native Client' (or OLE DB provider)\n"
     "Specify server name, authentication, and target database"),
    (ORANGE,   "4. Map Columns",
     "Select destination table (existing or create new)\n"
     "Review column mappings — check data types match source\n"
     "Enable 'Enable identity insert' if loading rows with explicit IDs"),
    (STEEL,    "5. Run or Save as SSIS Package",
     "Run immediately — or save as .dtsx SSIS package for reuse\n"
     "Review execution results — note rows transferred and any errors"),
    (AMBER,    "6. Validate the Import",
     "SELECT COUNT(*) FROM <table>  — row count matches source file?\n"
     "SELECT TOP 5 * FROM <table>  — spot-check data quality"),
]
for i, (col, title, desc) in enumerate(steps_wizard):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.62) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(0.38), Inches(1.45), col)
    tb(s, title, x + Inches(0.55), y + Inches(0.08), Inches(5.55), Inches(0.32),
       fs=14, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.55), y + Inches(0.48), Inches(5.55), Inches(0.82),
       fs=12, color=MID_GRAY)
footer(s, 3)


# ─── Slide 4: BULK INSERT ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "BULK INSERT", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "BULK INSERT loads a flat file directly into a table — faster than INSERT row by row and fully scriptable.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

bulk_sql = (
    "-- Create staging table for product import\n"
    "CREATE TABLE dbo.ProductImport (\n"
    "  ProductCode  VARCHAR(20),\n"
    "  ProductName  NVARCHAR(100),\n"
    "  Category     NVARCHAR(50),\n"
    "  UnitPrice    DECIMAL(10,2),\n"
    "  StockQty     INT\n"
    ");\n\n"
    "-- Load CSV file (comma-delimited, header row)\n"
    "BULK INSERT dbo.ProductImport\n"
    "FROM 'C:\\Data\\products.csv'\n"
    "WITH (\n"
    "  FIRSTROW = 2,          -- skip header row\n"
    "  FIELDTERMINATOR = ',', -- comma delimiter\n"
    "  ROWTERMINATOR = '\\n', -- newline row separator\n"
    "  TABLOCK                -- table-level lock for speed\n"
    ");\n\n"
    "-- Verify\n"
    "SELECT COUNT(*) AS LoadedRows FROM dbo.ProductImport;"
)
code(s, bulk_sql, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.55), fs=12)

annotations = [
    (ORANGE,   "FIRSTROW = 2",        "Skip the header row — FIRSTROW = 1 if no header."),
    (BLUE_ACC, "FIELDTERMINATOR",     "Match the delimiter in your CSV: ',' or '\\t' or '|'."),
    (GREEN_OK, "ROWTERMINATOR",       "Usually '\\n' (Unix) or '\\r\\n' (Windows CRLF)."),
    (STEEL,    "TABLOCK",             "Acquires table lock — improves bulk load throughput."),
    (AMBER,    "BATCHSIZE",           "BATCHSIZE = 10000 commits every 10K rows — safer for large files."),
    (NAVY,     "ERRORFILE",           "ERRORFILE = 'C:\\errors.txt' — logs bad rows to a file."),
    (ORANGE,   "Permissions required","Requires INSERT on table + ADMINISTER BULK OPERATIONS server permission."),
]
for i, (col, kw, note) in enumerate(annotations):
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.52), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(5.0), Inches(0.25),
       fs=13, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.30), Inches(5.0), Inches(0.28),
       fs=12, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: BCP Utility ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "BCP — Bulk Copy Program", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "BCP is a command-line tool for high-speed import and export between SQL Server and flat files.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

bcp_commands = (
    "-- Export a table to CSV\n"
    "bcp NovaMart.dbo.Employee out C:\\Data\\employees.csv\n"
    "  -S localhost -T -c -t,\n\n"
    "-- Import from CSV into a table\n"
    "bcp NovaMart.dbo.ProductImport in C:\\Data\\products.csv\n"
    "  -S localhost -T -c -t, -F 2\n\n"
    "-- Export a query result\n"
    "bcp \"SELECT * FROM NovaMart.dbo.Employee WHERE DeptID = 1\"\n"
    "  queryout C:\\Data\\dept1.csv\n"
    "  -S localhost -T -c -t,\n\n"
    "-- Key switches:\n"
    "-- -S  server name\n"
    "-- -T  Windows trusted connection  (-U login -P password for SQL auth)\n"
    "-- -c  character mode (text)\n"
    "-- -t  field delimiter\n"
    "-- -F  first row (skip header)"
)
code(s, bcp_commands, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.55), fs=12)

bcp_notes = [
    (STEEL,    "out / in / queryout",  "out = export table, in = import, queryout = export query result."),
    (BLUE_ACC, "-T (trusted)",         "Uses Windows auth; replace with -U <user> -P <pass> for SQL auth."),
    (GREEN_OK, "-c character mode",    "Most compatible — outputs plain text. Use -n for native binary."),
    (ORANGE,   "-t field delimiter",   "Matches BULK INSERT FIELDTERMINATOR — keep them consistent."),
    (AMBER,    "-F first row",         "-F 2 skips the header line in the input file."),
    (NAVY,     "Format file",          "bcp ... format nul -c -f schema.fmt — reusable column mapping."),
    (STEEL,    "BCP vs BULK INSERT",   "BCP = command prompt, no SQL context. BULK INSERT = T-SQL, more control."),
]
for i, (col, kw, note) in enumerate(bcp_notes):
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.52), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(5.0), Inches(0.25),
       fs=13, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.30), Inches(5.0), Inches(0.28),
       fs=12, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: ETL Concepts ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "ETL Concepts & SSIS Introduction", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "ETL (Extract, Transform, Load) is the pattern behind every serious data integration workflow.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

etl_stages = [
    (BLUE_ACC, "EXTRACT",
     "Read data from source systems.\n"
     "Sources: CSV files, Excel sheets, other SQL Server databases,\n"
     "REST APIs, ERP systems.\n"
     "Key concern: minimal impact on source during extraction."),
    (ORANGE,   "TRANSFORM",
     "Clean, validate, and reshape data.\n"
     "Examples: convert date formats, standardise codes,\n"
     "split full name into first/last, filter invalid rows,\n"
     "apply lookup tables (e.g. DeptCode → DeptID)."),
    (GREEN_OK, "LOAD",
     "Write transformed data to the destination.\n"
     "Strategies: full replace (truncate + load), incremental\n"
     "(insert new rows only), upsert (MERGE statement).\n"
     "Key concern: performance, locks, and rollback safety."),
]
for i, (col, title, desc) in enumerate(etl_stages):
    x = Inches(0.45) + i * Inches(4.28)
    rect(s, x, Inches(1.6), Inches(3.98), Inches(4.55), col)
    tb(s, title, x, Inches(1.72), Inches(3.98), Inches(0.45),
       fs=26, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc, x + Inches(0.2), Inches(2.32), Inches(3.6), Inches(3.5),
       fs=14, color=WHITE)

tb(s, "SSIS (SQL Server Integration Services) — Visual Studio-based tool to build ETL packages with drag-and-drop data flows.",
   Inches(0.45), Inches(6.28), Inches(12.45), Inches(0.38), fs=14, color=DARK_TEXT)
merge_sql = (
    "MERGE dbo.Product AS target\n"
    "USING dbo.ProductImport AS source ON target.ProductCode = source.ProductCode\n"
    "WHEN MATCHED THEN UPDATE SET target.UnitPrice = source.UnitPrice\n"
    "WHEN NOT MATCHED THEN INSERT (ProductCode, ProductName, UnitPrice)\n"
    "  VALUES (source.ProductCode, source.ProductName, source.UnitPrice);"
)
code(s, merge_sql, Inches(0.45), Inches(6.72), Inches(12.45), Inches(0.62), fs=11)
footer(s, 6)


# ─── Slide 7: Lab 11 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), GREEN_OK)
tb(s, "LAB 11", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Migrate NovaMart's Legacy Product Catalog", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~40 min  ·  Tool: SSMS + Import Wizard  ·  Database: NovaMart",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Create products.csv on disk with columns: ProductCode, ProductName, Category, UnitPrice, StockQty",
    "Add 5-10 sample product rows (e.g. NM-001, Running Shoes, Footwear, 299000, 50)",
    "Create staging table: CREATE TABLE dbo.ProductImport (...) in NovaMart",
    "Use SSMS Import Wizard (Tasks → Import Data) to load products.csv into ProductImport",
    "Validate: SELECT COUNT(*) — confirms row count",
    "Use BULK INSERT to reload the same file into ProductImport (truncate first)",
    "Export Employee table to CSV: Tasks → Export Data → Flat File destination",
    "Bonus: write a MERGE statement to upsert ProductImport into a Product table",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.47) * i, Inches(0.32), Inches(0.32), GREEN_OK)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.47) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.47) * i, Inches(6.8), Inches(0.36),
       fs=12, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "products.csv file created on disk",
    "Wizard import completes — row count matches",
    "BULK INSERT reloads without error",
    "Employee export CSV file created",
    "Bonus: MERGE runs without conflicts",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.70) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)
footer(s, 7)


# ─── Slide 8: Course Wrap-Up ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Course Complete", Inches(0), Inches(0.45), SLIDE_W, Inches(0.65),
   fs=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "NovaMart is now running — backed up, monitored, secured, and optimised.",
   Inches(0), Inches(1.18), SLIDE_W, Inches(0.42), fs=20, color=ORANGE, align=PP_ALIGN.CENTER)

outcomes_all = [
    (STEEL,    "Install & Configure",  "SQL Server installed, TCP/IP enabled, services configured"),
    (BLUE_ACC, "Database Management",  "NovaMart created with correct files, filegroups, and recovery model"),
    (GREEN_OK, "Tables & Data",        "Department and Employee tables with constraints, CRUD operations"),
    (ORANGE,   "Security",             "SQL logins, database users, roles, GRANT/DENY/REVOKE tested"),
    (STEEL,    "Backup & Restore",     "Full + Diff + Log backups, disaster recovery simulation completed"),
    (BLUE_ACC, "Maintenance",          "Index rebuild, statistics update, DBCC CHECKDB — all green"),
    (GREEN_OK, "Automation",           "Nightly backup job scheduled and tested in SQL Server Agent"),
    (ORANGE,   "Monitoring",           "Activity Monitor, DMVs, blocking chain identified and resolved"),
    (STEEL,    "Troubleshooting",      "Connection errors, login failures, slow queries diagnosed"),
    (GREEN_OK, "Import / Export",      "CSV imported via Wizard and BULK INSERT, table exported to CSV"),
]
for i, (col, title, desc) in enumerate(outcomes_all):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.65) + row_idx * Inches(1.02)
    rect(s, x, y, Inches(0.28), Inches(0.85), col)
    tb(s, title, x + Inches(0.42), y + Inches(0.05), Inches(2.6), Inches(0.28),
       fs=13, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.42), y + Inches(0.42), Inches(5.5), Inches(0.32),
       fs=11, color=MID_GRAY)

tb(s, "Thank you — questions welcome. Lab guide and slide deck available on the course portal.",
   Inches(0), SLIDE_H - Inches(0.55), SLIDE_W, Inches(0.35),
   fs=14, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-11-Data-Import-and-Export.pptx")
prs.save(out)
print(f"Saved: {out}")
