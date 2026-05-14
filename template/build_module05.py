from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "05"
MOD_TITLE = "Security & User Management"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), STEEL)
tb(s, "MODULE 05", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Security &\nUser Management", Inches(0.55), Inches(1.65),
   Inches(12.2), Inches(1.6), fs=46, bold=True, color=WHITE)
tb(s, "Create SQL Server logins, map database users, assign roles,\nand verify access control — the foundation of secure administration.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=20, color=BLUE_ACC)
chips = [
    ("Logins vs Users",     ORANGE,   Inches(0.55)),
    ("Server & DB Roles",   GREEN_OK, Inches(4.45)),
    ("Access Control",      BLUE_ACC, Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(5.1), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.17), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), STEEL)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 05 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (STEEL,    "Explain the difference between a SQL Server login and a database user"),
    (ORANGE,   "Create a SQL Server login using both SSMS and T-SQL"),
    (BLUE_ACC, "Map a database user to a login and grant access to a specific database"),
    (GREEN_OK, "Assign server-level and database-level roles to control permissions"),
    (STEEL,    "Test and verify access control by connecting as the new login"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Logins vs Users ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Logins vs Database Users", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "These are two separate objects at two different levels — a common source of confusion.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Login (server-level) box
rect(s, Inches(0.45), Inches(1.55), Inches(5.75), Inches(5.0), NAVY)
tb(s, "LOGIN", Inches(0.45), Inches(1.72), Inches(5.75), Inches(0.45),
   fs=24, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
tb(s, "Server-Level Principal", Inches(0.45), Inches(2.22), Inches(5.75), Inches(0.3),
   fs=15, color=MID_GRAY, align=PP_ALIGN.CENTER)
login_points = [
    "Authenticates to the SQL Server instance",
    "Lives in master database (sys.server_principals)",
    "Can be Windows account, Windows group, or SQL login",
    "Has server-level role memberships (e.g., sysadmin)",
    "Created with: CREATE LOGIN ... WITH PASSWORD",
    "Does NOT automatically give database access",
]
for j, pt in enumerate(login_points):
    tb(s, f"·  {pt}", Inches(0.65), Inches(2.72) + j * Inches(0.52),
       Inches(5.35), Inches(0.4), fs=14, color=WHITE)

# Arrow
tb(s, "MAPS TO  →", Inches(6.3), Inches(3.72), Inches(1.2), Inches(0.38),
   fs=16, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

# User (database-level) box
rect(s, Inches(7.55), Inches(1.55), Inches(5.33), Inches(5.0), STEEL)
tb(s, "DATABASE USER", Inches(7.55), Inches(1.72), Inches(5.33), Inches(0.45),
   fs=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Database-Level Principal", Inches(7.55), Inches(2.22), Inches(5.33), Inches(0.3),
   fs=15, color=RGBColor(0xC5, 0xD8, 0xF0), align=PP_ALIGN.CENTER)
user_points = [
    "Authorises access inside one database",
    "Lives in the database (sys.database_principals)",
    "Linked to a login via SID (security identifier)",
    "Has database-level role memberships",
    "Created with: CREATE USER ... FOR LOGIN",
    "Must exist in each database the login needs",
]
for j, pt in enumerate(user_points):
    tb(s, f"·  {pt}", Inches(7.75), Inches(2.72) + j * Inches(0.52),
       Inches(5.0), Inches(0.4), fs=14, color=WHITE)

tb(s, "One login can be mapped to users in many databases. One database user corresponds to exactly one login.",
   Inches(0.45), Inches(6.65), Inches(12.45), Inches(0.32), fs=13, italic=True, color=MID_GRAY)
footer(s, 3)


# ─── Slide 4: Server Roles & Database Roles ───────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Server Roles & Database Roles", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Roles bundle permissions together — assign roles rather than individual object permissions wherever possible.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

server_roles = [
    ("sysadmin",      "Full unrestricted access to the instance. SA login is a member."),
    ("serveradmin",   "Can configure server-wide settings (sp_configure, SHUTDOWN)."),
    ("securityadmin", "Can manage logins and grant/revoke server permissions."),
    ("dbcreator",     "Can create, alter, drop, and restore any database."),
    ("bulkadmin",     "Can run BULK INSERT statements."),
]
db_roles = [
    ("db_owner",          "Full control over the database — assign sparingly."),
    ("db_datareader",     "SELECT on all tables and views in the database."),
    ("db_datawriter",     "INSERT, UPDATE, DELETE on all tables."),
    ("db_ddladmin",       "Can run DDL (CREATE TABLE, ALTER, DROP) — no data read."),
    ("db_securityadmin",  "Manages database-level role membership and permissions."),
    ("public",            "Baseline role every user is a member of automatically."),
]

rect(s, Inches(0.45), Inches(1.62), Inches(5.95), Inches(4.75), NAVY)
tb(s, "SERVER ROLES  (instance-level)", Inches(0.45), Inches(1.78), Inches(5.95), Inches(0.38),
   fs=16, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
for j, (role, desc) in enumerate(server_roles):
    y = Inches(2.32) + j * Inches(0.72)
    tb(s, role, Inches(0.65), y, Inches(1.75), Inches(0.3), fs=13, bold=True, color=WHITE)
    tb(s, desc, Inches(2.45), y, Inches(3.78), Inches(0.42), fs=12,
       color=RGBColor(0xC5, 0xD8, 0xF0))

rect(s, Inches(6.75), Inches(1.62), Inches(6.13), Inches(4.75), STEEL)
tb(s, "DATABASE ROLES  (per database)", Inches(6.75), Inches(1.78), Inches(6.13), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
for j, (role, desc) in enumerate(db_roles):
    y = Inches(2.32) + j * Inches(0.58)
    tb(s, role, Inches(6.95), y, Inches(2.0), Inches(0.28), fs=12, bold=True, color=WHITE)
    tb(s, desc, Inches(9.0), y, Inches(3.7), Inches(0.35), fs=11,
       color=RGBColor(0xC5, 0xD8, 0xF0))
footer(s, 4)


# ─── Slide 5: Creating Logins & Users via T-SQL ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "Creating Logins & Users via T-SQL", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "T-SQL is the repeatable, scriptable way to manage SQL Server security principals.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

create_sql = (
    "-- 1. Create a SQL Server login (instance level)\n"
    "CREATE LOGIN AppUser\n"
    "WITH PASSWORD = 'Str0ng!Passw0rd',\n"
    "     CHECK_POLICY = ON,\n"       # enforces Windows password policy
    "     DEFAULT_DATABASE = LabDB_TSQL;\n\n"
    "-- 2. Create a database user mapped to the login\n"
    "USE LabDB_TSQL;\n"
    "GO\n"
    "CREATE USER AppUser FOR LOGIN AppUser;\n\n"
    "-- 3. Assign a database role\n"
    "ALTER ROLE db_datareader ADD MEMBER AppUser;\n\n"
    "-- 4. Grant specific permission (optional)\n"
    "GRANT SELECT ON dbo.Employee TO AppUser;\n"
    "GRANT EXECUTE ON dbo.usp_GetSalary TO AppUser;"
)
code(s, create_sql, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.1), fs=12)

annotations = [
    (NAVY,     "CHECK_POLICY = ON",     "Enforces Windows password complexity rules"),
    (BLUE_ACC, "DEFAULT_DATABASE",       "Database to connect to if none specified"),
    (GREEN_OK, "FOR LOGIN",              "Links this database user to the server login"),
    (ORANGE,   "ALTER ROLE ... ADD MEMBER", "Grants all permissions of the role to the user"),
    (STEEL,    "GRANT SELECT ON <table>",  "Fine-grained permission on a specific object"),
    (NAVY,     "GRANT EXECUTE ON <proc>",  "Permission to run a stored procedure"),
    (BLUE_ACC, "REVOKE / DENY",            "REVOKE removes a grant; DENY overrides all"),
]
for i, (col, kw, note) in enumerate(annotations):
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.52), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(5.0), Inches(0.25),
       fs=13, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.30), Inches(5.0), Inches(0.28),
       fs=12, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: SSMS Security Workflow ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Managing Security via SSMS", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "The SSMS GUI mirrors T-SQL operations — useful for learning; T-SQL is preferred for production scripting.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
workflow = [
    (NAVY,     "Create a Login (SSMS)",
     "Instance → Security → Logins → right-click → New Login\n"
     "Enter name, choose auth type, set default database."),
    (BLUE_ACC, "Map to a Database User",
     "In the New Login dialog → User Mapping page\n"
     "Check the target database → assign role membership (db_datareader etc.)."),
    (GREEN_OK, "Test the Login",
     "New connection in SSMS → specify the login credentials\n"
     "Attempt SELECT from a table — confirm access or permission error."),
    (ORANGE,   "Disable / Drop a Login",
     "Right-click Login → Properties → Status tab → set Login to Disabled\n"
     "Or: ALTER LOGIN AppUser DISABLE; — DROP LOGIN removes permanently."),
    (STEEL,    "Audit Permissions",
     "Object Explorer → Database → Security → Users → right-click → Properties\n"
     "Or query: sys.database_role_members, sys.server_role_members, fn_my_permissions."),
    (NAVY,     "sa Login Best Practice",
     "Rename sa (ALTER LOGIN sa WITH NAME = ...)\n"
     "Disable it entirely if Windows Auth is the only mode in use."),
]
for i, (col, title, desc) in enumerate(workflow):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.62) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(0.38), Inches(1.48), col)
    tb(s, title, x + Inches(0.55), y + Inches(0.08), Inches(5.55), Inches(0.32),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.55), y + Inches(0.48), Inches(5.55), Inches(0.75),
       fs=13, color=MID_GRAY)
footer(s, 6)


# ─── Slide 7: Lab 05 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), STEEL)
tb(s, "LAB 05", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Create a Login, Map a User & Test Access Control", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~45 min  ·  Tool: SSMS  ·  Uses: LabDB_TSQL from Module 04",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Create SQL login 'LabReader' via SSMS: Security → Logins → New Login (SQL Auth)",
    "Map LabReader to LabDB_TSQL with role db_datareader — via User Mapping page",
    "Open a second SSMS connection using LabReader credentials",
    "As LabReader: SELECT * FROM Employee — should succeed",
    "As LabReader: try INSERT INTO Employee … — confirm permission denied",
    "Back as sysadmin: ALTER ROLE db_datawriter ADD MEMBER LabReader",
    "As LabReader: retry INSERT — should now succeed",
    "Run: SELECT IS_MEMBER('db_datawriter') — confirm returns 1",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=STEEL)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.45) * i, Inches(0.32), Inches(0.32), STEEL)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.45) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.45) * i, Inches(6.8), Inches(0.36),
       fs=13, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "Login appears in Security → Logins",
    "SELECT succeeds as LabReader",
    "INSERT denied before role change",
    "INSERT succeeds after role change",
    "IS_MEMBER returns 1",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.70) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)
tb(s, "Clean up: DROP USER LabReader from LabDB_TSQL; DROP LOGIN LabReader.",
   Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3), fs=12, italic=True, color=MID_GRAY)
footer(s, 7)


# ─── Slide 8: Day 1 Wrap-Up ──────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Day 1 Complete", Inches(0), Inches(0.45), SLIDE_W, Inches(0.65),
   fs=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "You have built a working SQL Server environment from the ground up.",
   Inches(0), Inches(1.18), SLIDE_W, Inches(0.42), fs=21, color=ORANGE, align=PP_ALIGN.CENTER)

day1_summary = [
    (STEEL,    "M01 — Architecture",   "Instance, databases, system DBs, auth modes, SSMS"),
    (BLUE_ACC, "M02 — Installation",   "Planning, setup wizard, Config Manager, TCP/IP"),
    (GREEN_OK, "M03 — Databases",      "MDF/LDF files, CREATE DATABASE, filegroups, properties"),
    (ORANGE,   "M04 — Tables & Data",  "Data types, constraints, INSERT/UPDATE/DELETE/SELECT"),
    (STEEL,    "M05 — Security",       "Logins, users, server roles, database roles, access test"),
]
for i, (col, title, desc) in enumerate(day1_summary):
    col_idx = i % 2
    row_idx = i // 2
    if i == 4:
        x = Inches(3.35)
        y = Inches(5.32)
    else:
        x = Inches(0.45) + col_idx * Inches(6.38)
        y = Inches(1.85) + row_idx * Inches(1.58)
    rect(s, x, y, Inches(5.95), Inches(1.35), col)
    tb(s, title, x + Inches(0.2), y + Inches(0.1), Inches(5.55), Inches(0.35),
       fs=16, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.2), y + Inches(0.52), Inches(5.55), Inches(0.55),
       fs=13, color=WHITE)

tb(s, "Day 2 begins with Backup & Restore. Have your instance running and SSMS open.",
   Inches(0), Inches(7.0), SLIDE_W, Inches(0.3), fs=14, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-05-Security-and-User-Management.pptx")
prs.save(out)
print(f"Saved: {out}")
