from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
RED_ERR  = RGBColor(0xE0, 0x52, 0x52)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "10"
MOD_TITLE = "Basic Troubleshooting"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 7

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), ORANGE)
tb(s, "MODULE 10", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Basic Troubleshooting", Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.2),
   fs=48, bold=True, color=WHITE)
tb(s, "Monday morning tickets: connection errors, login failures, and slow queries.\n"
      "Here is how to diagnose and resolve them confidently.",
   Inches(0.55), Inches(3.35), Inches(11.0), Inches(0.9), fs=19, color=ORANGE)
chips = [
    ("Connection Issues", RED_ERR,  Inches(0.55)),
    ("Login Failures",    AMBER,    Inches(4.45)),
    ("Error Logs",        BLUE_ACC, Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(4.95), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.02), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Day 2  ·  Afternoon  ·  1 hour", Inches(0.55), Inches(6.72), Inches(6.0), Inches(0.35),
   fs=13, color=MID_GRAY)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 10 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (ORANGE,   "Diagnose and resolve the most common SQL Server connection errors"),
    (RED_ERR,  "Identify the root cause of login failure errors (18456 and variants)"),
    (BLUE_ACC, "Read the SQL Server Error Log in SSMS and via T-SQL"),
    (GREEN_OK, "Apply a structured first-response workflow to an unfamiliar incident"),
    (STEEL,    "Simulate and resolve three realistic Monday-morning support tickets"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Connection Errors ──────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), RED_ERR)
tb(s, "Common Connection Errors", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Most connection failures have a small set of root causes — check them in order.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

errors = [
    (RED_ERR,  'Error 53 — "Network path not found"',
     "SQL Server Browser not running, or TCP/IP disabled.\n"
     "Fix: SQL Server Configuration Manager → enable TCP/IP → restart service.\n"
     "Check: telnet <server> 1433  or  Test-NetConnection <server> -Port 1433"),
    (AMBER,    'Error 17 — "SQL Server does not exist or access denied"',
     "Instance name wrong, or firewall blocking port 1433.\n"
     "Fix: Verify instance name (SERVERNAME or SERVERNAME\\INSTANCENAME).\n"
     "Check: Windows Firewall → Inbound Rule for TCP 1433 exists and is enabled."),
    (ORANGE,   'Error 4060 — "Cannot open database requested"',
     "Login is valid but default database is offline or dropped.\n"
     "Fix: ALTER LOGIN [user] WITH DEFAULT_DATABASE = master;\n"
     "Or: specify a different database in the connection string."),
    (STEEL,    'Error 10060 — "Connection timed out"',
     "Network latency, wrong IP, or SQL Server not accepting connections.\n"
     "Fix: ping the server, check IP/hostname, verify SQL Server is Running.\n"
     "Check: sys.dm_exec_connections — if 0 rows, service may be down."),
]
for i, (col, title, desc) in enumerate(errors):
    y = Inches(1.58) + i * Inches(1.38)
    rect(s, Inches(0.45), y, Inches(0.38), Inches(1.15), col)
    tb(s, title, Inches(1.02), y + Inches(0.05), Inches(11.85), Inches(0.32),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, Inches(1.02), y + Inches(0.42), Inches(11.85), Inches(0.65),
       fs=13, color=MID_GRAY)
footer(s, 3)


# ─── Slide 4: Login Failures (Error 18456) ───────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), AMBER)
tb(s, 'Login Failure — Error 18456', Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, 'Error 18456 means "login failed." The state number reveals the exact reason.',
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

states = [
    ("State 1",  "Generic — details suppressed for security. Check Error Log for real state."),
    ("State 2",  "Invalid user ID — login not found in sys.server_principals."),
    ("State 5",  "Invalid user ID — same as state 2 (older versions)."),
    ("State 6",  "Login attempted with Windows auth but only SQL auth is configured."),
    ("State 7",  "Login disabled AND wrong password provided."),
    ("State 8",  "Correct login name but wrong password."),
    ("State 9",  "Login does not have permission to log in — login disabled."),
    ("State 11", "Login valid but access denied to the server — check server role."),
    ("State 18", "Password expired — must change at next login."),
    ("State 38", "Login valid, but database is unavailable (offline, restoring, or dropped)."),
]
rect(s, Inches(0.45), Inches(1.55), Inches(12.45), Inches(4.65), NAVY)
tb(s, "STATE", Inches(0.65), Inches(1.68), Inches(1.4), Inches(0.32),
   fs=14, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
tb(s, "MEANING", Inches(2.22), Inches(1.68), Inches(10.5), Inches(0.32),
   fs=14, bold=True, color=AMBER)
for j, (state, meaning) in enumerate(states):
    y = Inches(2.08) + j * Inches(0.42)
    bg_col = NAVY if j % 2 == 0 else RGBColor(0x10, 0x22, 0x4E)
    rect(s, Inches(0.45), y, Inches(12.45), Inches(0.42), bg_col)
    tb(s, state, Inches(0.65), y + Inches(0.06), Inches(1.4), Inches(0.28),
       fs=12, bold=True, color=AMBER, align=PP_ALIGN.CENTER)
    tb(s, meaning, Inches(2.22), y + Inches(0.06), Inches(10.5), Inches(0.28),
       fs=12, color=WHITE)

tb(s, "Always check the SQL Server Error Log for the actual state — client messages may show only State 1.",
   Inches(0.45), Inches(6.28), Inches(12.45), Inches(0.32), fs=13, italic=True, color=DARK_TEXT)
read_log_sql = "EXEC sp_readerrorlog 0, 1, 'Login failed', 'NovaMart';"
code(s, read_log_sql, Inches(0.45), Inches(6.65), Inches(8.0), Inches(0.40), fs=12)
footer(s, 4)


# ─── Slide 5: Reading the SQL Server Error Log ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "Reading the SQL Server Error Log", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "The Error Log is the authoritative record of everything SQL Server has experienced.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Two-panel layout
rect(s, Inches(0.45), Inches(1.55), Inches(5.85), Inches(4.55), NAVY)
tb(s, "VIA SSMS", Inches(0.45), Inches(1.68), Inches(5.85), Inches(0.35),
   fs=16, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
ssms_steps = [
    "In Object Explorer: Management → SQL Server Logs",
    "Current log = most recent; Archive #1 = previous, etc.",
    "Double-click 'Current' to open Log File Viewer",
    "Filter by date range, error level, or search text",
    "Look for: Severity 16+ errors, stack dumps, I/O errors",
    "Event types: Information / Warning / Error",
    "Logs cycle at restart or when sp_cycle_errorlog is run",
]
for j, step in enumerate(ssms_steps):
    tb(s, f"·  {step}", Inches(0.65), Inches(2.18) + j * Inches(0.54),
       Inches(5.45), Inches(0.42), fs=13, color=WHITE)

rect(s, Inches(6.75), Inches(1.55), Inches(6.13), Inches(4.55), RGBColor(0x05, 0x12, 0x2A))
tb(s, "VIA T-SQL", Inches(6.75), Inches(1.68), Inches(6.13), Inches(0.35),
   fs=16, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
log_sql = (
    "-- Read current error log\n"
    "EXEC sp_readerrorlog;\n\n"
    "-- Search for a keyword\n"
    "EXEC sp_readerrorlog 0, 1, 'error';\n\n"
    "-- Read archive log #1\n"
    "EXEC sp_readerrorlog 1, 1, 'Login failed';\n\n"
    "-- Check log file location\n"
    "EXEC xp_readerrorlog 0, 1, N'Logging SQL Server messages';\n\n"
    "-- Cycle the log (keeps current small)\n"
    "EXEC sp_cycle_errorlog;"
)
code(s, log_sql, Inches(6.88), Inches(2.05), Inches(5.88), Inches(3.88), fs=12)

tb(s, "Critical errors to act on immediately:  823 (I/O failure), 824 (page checksum fail), 9001 (log not available).",
   Inches(0.45), Inches(6.22), Inches(12.45), Inches(0.35), fs=13, bold=True, color=DARK_TEXT)
footer(s, 5)


# ─── Slide 6: Performance Bottleneck Basics ──────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "Performance Bottlenecks — First Response", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Narrow down to the right layer before tuning anything.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

bottlenecks = [
    (ORANGE,   "CPU Bottleneck",
     "Symptoms: high CPU in Task Manager, high Batch Requests/sec, CXPACKET waits.\n"
     "Quick checks: sp_who2, dm_exec_requests ORDER BY cpu_time DESC.\n"
     "Common causes: missing indexes → table scans, implicit conversions, high MAXDOP."),
    (BLUE_ACC, "Memory Bottleneck",
     "Symptoms: high PAGEIOLATCH waits, low PLE (Page Life Expectancy < 300 sec).\n"
     "Quick checks: dm_os_sys_info, dm_os_buffer_descriptors.\n"
     "Causes: max server memory set too low, large unoptimised queries."),
    (RED_ERR,  "I/O Bottleneck",
     "Symptoms: PAGEIOLATCH_SH/_EX waits, slow reads in Data File I/O chart.\n"
     "Quick checks: dm_io_virtual_file_stats, DBCC SQLPERF('sys.dm_os_wait_stats').\n"
     "Causes: storage latency, log drive contention, no SSD for log file."),
    (GREEN_OK, "Lock / Blocking Bottleneck",
     "Symptoms: LCK_M_X waits, sessions stuck on 'SUSPENDED' status.\n"
     "Quick checks: dm_exec_requests WHERE blocking_session_id > 0.\n"
     "Causes: long open transactions, missing indexes causing wide table locks."),
    (STEEL,    "Network Bottleneck",
     "Symptoms: ASYNC_NETWORK_IO waits, application-side latency despite fast queries.\n"
     "Quick checks: dm_exec_sessions — check reads/writes vs elapsed time.\n"
     "Causes: fetching large result sets, client-side slow row processing."),
    (AMBER,    "Query Design Bottleneck",
     "Symptoms: high logical reads, Table Scan operators in execution plan.\n"
     "Quick checks: dm_exec_query_stats ORDER BY total_logical_reads DESC.\n"
     "Causes: missing covering indexes, SELECT *, outdated statistics."),
]
for i, (col, title, desc) in enumerate(bottlenecks):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.62) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(0.38), Inches(1.48), col)
    tb(s, title, x + Inches(0.55), y + Inches(0.08), Inches(5.55), Inches(0.32),
       fs=14, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.55), y + Inches(0.48), Inches(5.55), Inches(0.82),
       fs=12, color=MID_GRAY)
footer(s, 6)


# ─── Slide 7: Lab 10 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), ORANGE)
tb(s, "LAB 10", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Monday Morning — 3 Support Tickets", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~40 min  ·  Tool: SSMS + Configuration Manager  ·  Instance: local",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

tickets = [
    (RED_ERR,  "Ticket 1: 'Cannot connect to SQL Server'",
     "Simulate: stop SQL Server Browser service → attempt named instance connection\n"
     "Diagnose: check service status in Config Manager\n"
     "Resolve: restart SQL Server Browser → reconnect"),
    (AMBER,    "Ticket 2: 'Login failed for user LabSari'",
     "Simulate: disable login LabSari (ALTER LOGIN LabSari DISABLE)\n"
     "Diagnose: read Error Log → identify state number\n"
     "Resolve: ALTER LOGIN LabSari ENABLE → verify login works"),
    (BLUE_ACC, "Ticket 3: 'The report query is very slow today'",
     "Simulate: run SELECT * FROM Employee without WHERE clause (large result)\n"
     "Diagnose: Activity Monitor → Recent Expensive Queries → view execution plan\n"
     "Resolve: add WHERE clause or proper filter → compare elapsed time"),
]
for i, (col, title, desc) in enumerate(tickets):
    y = Inches(2.62) + i * Inches(1.38)
    rect(s, Inches(0.5), y, Inches(0.38), Inches(1.18), col)
    tb(s, title, Inches(1.05), y + Inches(0.08), Inches(11.75), Inches(0.32),
       fs=15, bold=True, color=WHITE)
    tb(s, desc, Inches(1.05), y + Inches(0.46), Inches(11.75), Inches(0.60),
       fs=13, color=MID_GRAY)

tb(s, "Validation: all 3 tickets resolved, root cause identified, and steps documented in a comment block in SSMS.",
   Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.28), fs=12, italic=True, color=MID_GRAY)
footer(s, 7)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-10-Basic-Troubleshooting.pptx")
prs.save(out)
print(f"Saved: {out}")
