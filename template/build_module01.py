from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "01"
MOD_TITLE = "SQL Server Overview & Architecture"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 7

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), ORANGE)
tb(s, "MODULE 01", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "SQL Server Overview\n& Architecture", Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.6),
   fs=46, bold=True, color=WHITE)
tb(s, "Understand the engine, editions, and authentication foundation\nthat every SQL Server administration task depends on.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=21, color=ORANGE)
chips = [
    ("Editions & Components", BLUE_ACC, Inches(0.55)),
    ("Instance Architecture",  GREEN_OK, Inches(4.45)),
    ("System Databases",       STEEL,    Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(5.1), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.17), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 01 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38),
   fs=18, italic=True, color=MID_GRAY)
objectives = [
    (ORANGE,   "Identify SQL Server editions and choose the right one for a given scenario"),
    (BLUE_ACC, "Describe the instance, database, and object hierarchy inside the engine"),
    (GREEN_OK, "Name the four system databases and explain their purpose"),
    (STEEL,    "Distinguish Windows Authentication from SQL Server Authentication"),
    (ORANGE,   "Connect to a SQL Server instance using SSMS and navigate the object tree"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52),
       fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: SQL Server Editions ─────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "SQL Server Editions", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Choose the edition that matches the workload — you can upgrade, but plan early.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
editions = [
    (GREEN_OK, "EXPRESS",    "Free",  ["Up to 10 GB per database",
                                        "Great for dev/test and small apps",
                                        "No SQL Server Agent"]),
    (BLUE_ACC, "STANDARD",   "$$$",   ["Core BI features included",
                                        "Max 24 cores / 128 GB RAM",
                                        "Most production workloads"]),
    (ORANGE,   "ENTERPRISE", "$$$$",  ["Unlimited cores and memory",
                                        "Always On, partitioning, HADR",
                                        "Full feature set"]),
    (NAVY,     "DEVELOPER",  "Free",  ["Feature-identical to Enterprise",
                                        "Dev and test only — not production",
                                        "Best choice for this course"]),
]
for i, (col, name, price, bullets) in enumerate(editions):
    x = Inches(0.45) + i * Inches(3.15)
    rect(s, x, Inches(1.72), Inches(2.9), Inches(4.7), col)
    tb(s, name,  x, Inches(1.92), Inches(2.9), Inches(0.42),
       fs=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, price, x, Inches(2.42), Inches(2.9), Inches(0.32),
       fs=14, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.15), Inches(2.88), Inches(2.6), Inches(0.02),
         RGBColor(0xFF, 0xFF, 0xFF))
    for j, b in enumerate(bullets):
        tb(s, f"· {b}", x + Inches(0.2), Inches(3.08) + Inches(0.62) * j,
           Inches(2.55), Inches(0.45), fs=14, color=WHITE)
footer(s, 3)


# ─── Slide 4: Architecture — Instance, Database, Objects ─────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), NAVY)
tb(s, "Instance, Database & Engine Architecture", Inches(0), Inches(0.38),
   SLIDE_W, Inches(0.62), fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "SQL Server is a layered system — understanding the layers makes every admin task clearer.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Three hierarchy layers
layers = [
    (NAVY,     "SQL SERVER INSTANCE",
     "A running service (sqlservr.exe). One server can host multiple named instances.\nDefault: MSSQLSERVER  |  Named: SERVER\\SQLEXPRESS"),
    (BLUE_ACC, "DATABASES",
     "Each instance holds system databases (master, model, msdb, tempdb) plus your user databases.\nData lives in .mdf files; transaction log in .ldf files."),
    (GREEN_OK, "OBJECTS",
     "Inside every database: tables, views, stored procedures, functions, triggers, indexes, users.\nManaged through SSMS Object Explorer or T-SQL."),
]
for i, (col, title, desc) in enumerate(layers):
    y = Inches(1.72) + i * Inches(1.62)
    rect(s, Inches(0.5), y, Inches(12.3), Inches(1.42), col)
    tb(s, title, Inches(0.7), y + Inches(0.12), Inches(3.5), Inches(0.38),
       fs=18, bold=True, color=WHITE)
    tb(s, desc, Inches(4.4), y + Inches(0.08), Inches(8.2), Inches(0.65),
       fs=15, color=WHITE)
    if i < 2:
        tb(s, "▼", Inches(6.3), y + Inches(1.42), Inches(1.0), Inches(0.25),
           fs=18, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# System DB reference
tb(s, "System databases:  master — server config    model — DB template    msdb — Agent & jobs    tempdb — working area",
   Inches(0.5), Inches(6.68), Inches(12.3), Inches(0.35), fs=13, italic=True, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: Authentication Modes ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Authentication Modes", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Choose during installation — changing later requires a service restart.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Left panel: Windows Auth
rect(s, Inches(0.5), Inches(1.6), Inches(5.8), Inches(5.1), NAVY)
tb(s, "WINDOWS\nAUTHENTICATION", Inches(0.5), Inches(1.78), Inches(5.8), Inches(0.72),
   fs=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
win_points = [
    "Uses Windows / Active Directory accounts",
    "No password stored inside SQL Server",
    "Integrated single sign-on (SSO)",
    "Kerberos ticket validation",
    "RECOMMENDED for domain environments",
]
for j, pt in enumerate(win_points):
    tb(s, f"✓  {pt}", Inches(0.75), Inches(2.72) + Inches(0.62) * j,
       Inches(5.1), Inches(0.45), fs=15, color=WHITE)

# Right panel: SQL Auth
rect(s, Inches(7.0), Inches(1.6), Inches(5.8), Inches(5.1), ORANGE)
tb(s, "SQL SERVER\nAUTHENTICATION", Inches(7.0), Inches(1.78), Inches(5.8), Inches(0.72),
   fs=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
sql_points = [
    "SQL login with username + password",
    "Password hash stored in master database",
    "Needed for non-Windows clients",
    "Required for some apps and cloud links",
    "Use Mixed Mode to allow both types",
]
for j, pt in enumerate(sql_points):
    tb(s, f"→  {pt}", Inches(7.22), Inches(2.72) + Inches(0.62) * j,
       Inches(5.35), Inches(0.45), fs=15, color=WHITE)

# VS divider
rect(s, Inches(6.1), Inches(2.6), Inches(0.7), Inches(0.7), LIGHT_BG)
tb(s, "VS", Inches(6.1), Inches(2.65), Inches(0.7), Inches(0.5),
   fs=20, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)

tb(s, "Best practice: enable Windows Authentication only. Add SQL logins only when the application requires it.",
   Inches(0.5), Inches(6.78), Inches(12.3), Inches(0.3), fs=13, italic=True, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: SSMS Overview ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "SQL Server Management Studio (SSMS)", Inches(0), Inches(0.38),
   SLIDE_W, Inches(0.62), fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "SSMS is the primary GUI tool for all SQL Server administration tasks in this course.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
areas = [
    (NAVY,     "Object Explorer",  "Left-side tree: instances → databases → tables, views, users, jobs, etc."),
    (BLUE_ACC, "Query Editor",     "Write and run T-SQL. Colour-coded syntax, execution plans, results grid."),
    (GREEN_OK, "Activity Monitor", "Live session, waits, and I/O metrics — covered in Module 09."),
    (ORANGE,   "Object Properties","Right-click any object to view and modify properties, permissions, and settings."),
    (STEEL,    "SQL Server Agent", "Schedule jobs and manage alerts — covered in Module 08."),
    (NAVY,     "Reports",          "Standard reports for disk usage, index stats, and more under each database."),
]
for i, (col, title, desc) in enumerate(areas):
    col_num = i % 2
    row_num = i // 2
    x = Inches(0.5) + col_num * Inches(6.45)
    y = Inches(1.72) + row_num * Inches(1.42)
    rect(s, x, y, Inches(0.38), Inches(0.88), col)
    rect(s, x + Inches(0.38), y, Inches(5.8), Inches(0.88), WHITE)
    tb(s, title, x + Inches(0.52), y + Inches(0.08), Inches(5.4), Inches(0.3),
       fs=16, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.52), y + Inches(0.42), Inches(5.4), Inches(0.38),
       fs=13, color=MID_GRAY)
footer(s, 6)


# ─── Slide 7: Lab 01 ──────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), ORANGE)
tb(s, "LAB 01", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Explore SSMS & Connect to an Instance", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=32, bold=True, color=WHITE)
tb(s, "Duration: ~30 min  ·  Tool: SSMS  ·  Assumes: SQL Server Developer installed or training VM ready",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

# Steps + outcome in two columns
steps = [
    "Open SSMS and use Connect > Database Engine",
    "Connect to the local default instance (.) or the VM hostname",
    "Expand the Databases node — note the four system databases",
    "Open a New Query window and run:  SELECT @@VERSION",
    "In Object Explorer → Security → Logins: identify the SA login",
    "Check Server Properties → Security tab for authentication mode",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(3.95), RGBColor(0x15, 0x28, 0x55))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=ORANGE)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.55) * i, Inches(0.32), Inches(0.32), ORANGE)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.55) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.55) * i, Inches(6.8), Inches(0.38),
       fs=14, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(3.95), RGBColor(0x15, 0x28, 0x55))
tb(s, "EXPECTED OUTCOME", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "SSMS opens and shows the instance",
    "SELECT @@VERSION returns the SQL Server build",
    "You can see master, model, msdb, tempdb",
    "Authentication mode is visible in properties",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.72) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)

tb(s, "Clean up: no objects were created — nothing to remove.",
   Inches(0.5), Inches(6.72), Inches(12.3), Inches(0.3), fs=12, italic=True, color=MID_GRAY)
footer(s, 7)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-01-SQL-Server-Overview-and-Architecture.pptx")
prs.save(out)
print(f"Saved: {out}")
