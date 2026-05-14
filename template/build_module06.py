from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
ORANGE_D = RGBColor(0xC0, 0x3E, 0x0F)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
RED_ERR  = RGBColor(0xE0, 0x52, 0x52)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "06"
MOD_TITLE = "Backup & Restore"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), ORANGE)
tb(s, "MODULE 06", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Backup & Restore", Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.2),
   fs=52, bold=True, color=WHITE)
tb(s, "Protect NovaMart's data with full, differential, and log backups —\n"
      "and practice restoring before disaster strikes.",
   Inches(0.55), Inches(3.35), Inches(11.0), Inches(0.9), fs=20, color=ORANGE)
chips = [
    ("Recovery Models",   BLUE_ACC, Inches(0.55)),
    ("Backup Types",      GREEN_OK, Inches(4.45)),
    ("Restore Strategy",  STEEL,    Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(4.95), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.02), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Day 2  ·  Morning  ·  2 hours", Inches(0.55), Inches(6.72), Inches(6.0), Inches(0.35),
   fs=13, color=MID_GRAY)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), ORANGE)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 06 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (ORANGE,   "Explain the difference between Simple and Full recovery models"),
    (BLUE_ACC, "Perform a full database backup using SSMS and T-SQL"),
    (GREEN_OK, "Perform a differential backup to capture incremental changes"),
    (STEEL,    "Back up and restore the transaction log to a point in time"),
    (ORANGE,   "Restore the NovaMart database from a backup file and verify data"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Recovery Models ────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "Recovery Models", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "The recovery model controls how much of the transaction log is retained — and what you can restore.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

models = [
    (BLUE_ACC, "SIMPLE",
     "Log truncated automatically at each checkpoint.\n"
     "No log backups possible.\n"
     "Maximum data loss = changes since last full or diff backup.\n"
     "Good for: development, test, non-critical databases."),
    (ORANGE,   "FULL",
     "Full transaction log kept until a log backup is taken.\n"
     "Point-in-time restore is possible.\n"
     "Requires regular log backups to prevent log file growth.\n"
     "Required for: NovaMart production database."),
    (STEEL,    "BULK-LOGGED",
     "Minimally logs bulk operations (BULK INSERT, SELECT INTO).\n"
     "Reduces log growth during large data loads.\n"
     "Point-in-time restore not available during bulk period.\n"
     "Typically used as a temporary switch during large imports."),
]
for i, (col, title, desc) in enumerate(models):
    x = Inches(0.45) + i * Inches(4.28)
    rect(s, x, Inches(1.6), Inches(3.98), Inches(4.50), col)
    tb(s, title, x, Inches(1.72), Inches(3.98), Inches(0.45),
       fs=24, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, desc, x + Inches(0.2), Inches(2.32), Inches(3.6), Inches(3.45),
       fs=14, color=WHITE)

alter_sql = "ALTER DATABASE NovaMart SET RECOVERY FULL;"
code(s, alter_sql, Inches(0.45), Inches(6.22), Inches(8.0), Inches(0.62), fs=13)
tb(s, "NovaMart should always use FULL recovery.", Inches(8.6), Inches(6.28), Inches(4.3),
   Inches(0.42), fs=14, italic=True, color=MID_GRAY)
footer(s, 3)


# ─── Slide 4: Backup Types ───────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Backup Types", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Use a combination of backup types to balance storage cost against recovery time and data loss.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

backup_types = [
    (NAVY,     "FULL Backup",
     "Backs up the entire database — data pages + enough log to make it consistent.\n"
     "Self-contained: can restore from this one file alone.\n"
     "Typically run: weekly or nightly (small DBs)."),
    (BLUE_ACC, "DIFFERENTIAL Backup",
     "Backs up all pages changed since the last FULL backup.\n"
     "Grows larger until the next full backup resets the baseline.\n"
     "Typically run: nightly or every few hours."),
    (GREEN_OK, "TRANSACTION LOG Backup",
     "Backs up the active portion of the transaction log.\n"
     "Enables point-in-time restore (STOPAT).\n"
     "Truncates the log — required in FULL recovery to control log growth.\n"
     "Typically run: every 15–60 min in production."),
]
for i, (col, title, desc) in enumerate(backup_types):
    y = Inches(1.62) + i * Inches(1.72)
    rect(s, Inches(0.45), y, Inches(0.38), Inches(1.45), col)
    tb(s, title, Inches(1.05), y + Inches(0.08), Inches(4.5), Inches(0.35),
       fs=18, bold=True, color=DARK_TEXT)
    tb(s, desc, Inches(1.05), y + Inches(0.52), Inches(11.8), Inches(0.85),
       fs=14, color=MID_GRAY)

tb(s, "Strategy for NovaMart:  Full (Sunday) → Differential (Mon–Sat) → Log every 30 min",
   Inches(0.45), Inches(6.60), Inches(12.45), Inches(0.3), fs=13, bold=True,
   color=DARK_TEXT)
footer(s, 4)


# ─── Slide 5: Performing Backups – T-SQL ─────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Performing Backups via T-SQL", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "T-SQL backup commands are scriptable, schedulable, and the foundation for SQL Agent jobs.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

backup_sql = (
    "-- Full backup\n"
    "BACKUP DATABASE NovaMart\n"
    "TO DISK = 'C:\\SQLBackup\\NovaMart_Full.bak'\n"
    "WITH FORMAT, INIT,\n"
    "     NAME = 'NovaMart Full Backup',\n"
    "     COMPRESSION, STATS = 10;\n\n"
    "-- Differential backup\n"
    "BACKUP DATABASE NovaMart\n"
    "TO DISK = 'C:\\SQLBackup\\NovaMart_Diff.bak'\n"
    "WITH DIFFERENTIAL, INIT,\n"
    "     NAME = 'NovaMart Differential Backup',\n"
    "     COMPRESSION, STATS = 10;\n\n"
    "-- Transaction log backup\n"
    "BACKUP LOG NovaMart\n"
    "TO DISK = 'C:\\SQLBackup\\NovaMart_Log.bak'\n"
    "WITH INIT, NAME = 'NovaMart Log Backup', STATS = 10;"
)
code(s, backup_sql, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.55), fs=12)

annotations = [
    (ORANGE,   "FORMAT, INIT",        "Overwrite existing backup file. Remove INIT to append."),
    (BLUE_ACC, "COMPRESSION",         "Reduces backup file size — recommended for speed and storage."),
    (GREEN_OK, "STATS = 10",          "Print progress every 10% — useful for large databases."),
    (STEEL,    "WITH DIFFERENTIAL",   "Captures only pages changed since the last full backup."),
    (AMBER,    "BACKUP LOG",          "Only valid when recovery model is FULL or BULK-LOGGED."),
    (NAVY,     "Backup to network",   "Replace disk path with \\\\server\\share\\file.bak for UNC paths."),
    (ORANGE,   "VERIFY after backup", "RESTORE VERIFYONLY FROM DISK = '...' — check integrity."),
]
for i, (col, kw, note) in enumerate(annotations):
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.52), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(5.0), Inches(0.25),
       fs=13, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.30), Inches(5.0), Inches(0.28),
       fs=12, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: Restoring a Database ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "Restoring a Database", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "A backup is worthless unless you have practised the restore. Always test your backups.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

restore_sql = (
    "-- Step 1: Restore the full backup (WITH NORECOVERY keeps DB in restoring state)\n"
    "RESTORE DATABASE NovaMart\n"
    "FROM DISK = 'C:\\SQLBackup\\NovaMart_Full.bak'\n"
    "WITH NORECOVERY, STATS = 10;\n\n"
    "-- Step 2: Apply differential backup\n"
    "RESTORE DATABASE NovaMart\n"
    "FROM DISK = 'C:\\SQLBackup\\NovaMart_Diff.bak'\n"
    "WITH NORECOVERY, STATS = 10;\n\n"
    "-- Step 3: Apply log backup — bring online\n"
    "RESTORE LOG NovaMart\n"
    "FROM DISK = 'C:\\SQLBackup\\NovaMart_Log.bak'\n"
    "WITH RECOVERY;\n\n"
    "-- Optional: restore to a point in time\n"
    "-- WITH RECOVERY, STOPAT = '2026-05-08T09:30:00';"
)
code(s, restore_sql, Inches(0.45), Inches(1.52), Inches(6.55), Inches(5.55), fs=12)

notes = [
    (RED_ERR,  "NORECOVERY",        "Leaves DB in restoring state — required when more backups follow."),
    (GREEN_OK, "RECOVERY",          "Brings database online — use only on the FINAL restore step."),
    (BLUE_ACC, "STOPAT",            "Point-in-time recovery — only available with FULL recovery model."),
    (ORANGE,   "REPLACE option",    "Add WITH REPLACE to overwrite an existing database by same name."),
    (STEEL,    "MOVE option",       "WITH MOVE 'NovaMart' TO '...mdf', MOVE 'NovaMart_log' TO '...ldf'"),
    (AMBER,    "Tail-log backup",   "Before restoring a live DB, back up the tail of the active log first."),
    (NAVY,     "Verify before restore", "RESTORE HEADERONLY FROM DISK='...' to see backup sets on file."),
]
for i, (col, kw, note) in enumerate(notes):
    y = Inches(1.52) + i * Inches(0.71)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.52), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(5.0), Inches(0.25),
       fs=13, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.30), Inches(5.0), Inches(0.28),
       fs=12, color=MID_GRAY)
footer(s, 6)


# ─── Slide 7: Lab 06 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), ORANGE)
tb(s, "LAB 06", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "NovaMart Disaster Recovery Simulation", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~60 min  ·  Tool: SSMS  ·  Database: NovaMart",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Confirm NovaMart is in FULL recovery: SELECT name, recovery_model_desc FROM sys.databases",
    "Create backup folder: C:\\SQLBackup (or verify it exists)",
    "Run full backup T-SQL to NovaMart_Full.bak — note completion time",
    "Insert a test row: INSERT INTO Employee VALUES (...) — note the time",
    "Run differential backup to NovaMart_Diff.bak",
    "Insert another test row — note time",
    "Run log backup to NovaMart_Log.bak",
    "Simulate disaster: DROP TABLE Employee — confirm table gone",
    "Restore full → differential → log WITH RECOVERY",
    "Verify Employee table and all rows are back",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.45), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=ORANGE)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.38) * i, Inches(0.42), Inches(0.32), ORANGE)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.38) * i, Inches(0.42), Inches(0.28),
       fs=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.22), Inches(3.10) + Inches(0.38) * i, Inches(6.73), Inches(0.32),
       fs=12, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.45), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "Backup files created in C:\\SQLBackup",
    "DROP TABLE removes Employee",
    "Restore completes without error",
    "Employee table exists after restore",
    "All rows present including inserted rows",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.75) * i,
       Inches(3.88), Inches(0.55), fs=13, color=WHITE)
footer(s, 7)


# ─── Slide 8: Day 2 Progress Check ───────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Day 2 Progress", Inches(0), Inches(0.45), SLIDE_W, Inches(0.65),
   fs=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Module 06 complete — NovaMart data is now protected.",
   Inches(0), Inches(1.18), SLIDE_W, Inches(0.42), fs=21, color=ORANGE, align=PP_ALIGN.CENTER)

day2_progress = [
    (ORANGE,   "M06 ✓ Backup & Restore",   "Full, Diff, Log backups · Recovery models · Disaster recovery sim"),
    (STEEL,    "M07  Database Maintenance", "Index rebuild/reorganize · Statistics · DBCC CHECKDB"),
    (BLUE_ACC, "M08  SQL Server Agent",     "Jobs, schedules, alerts · Automate nightly backup"),
    (GREEN_OK, "M09  Monitoring & Perf.",   "Activity Monitor · DMVs · Blocking & deadlocks"),
    (MID_GRAY, "M10  Troubleshooting",      "Connection errors · Error logs · Performance bottlenecks"),
    (MID_GRAY, "M11  Import / Export",      "Wizard · CSV/Excel · BCP · ETL concepts"),
]
for i, (col, title, desc) in enumerate(day2_progress):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.85) + row_idx * Inches(1.52)
    rect(s, x, y, Inches(5.95), Inches(1.28), col)
    tb(s, title, x + Inches(0.2), y + Inches(0.1), Inches(5.55), Inches(0.35),
       fs=15, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.2), y + Inches(0.52), Inches(5.55), Inches(0.52),
       fs=12, color=WHITE)
tb(s, "Next up: Module 07 — Database Maintenance",
   Inches(0), Inches(7.0), SLIDE_W, Inches(0.3), fs=14, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-06-Backup-and-Restore.pptx")
prs.save(out)
print(f"Saved: {out}")
