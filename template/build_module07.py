from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
AMBER    = RGBColor(0xF5, 0x9E, 0x0B)
RED_ERR  = RGBColor(0xE0, 0x52, 0x52)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "07"
MOD_TITLE = "Database Maintenance"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 7

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), GREEN_OK)
tb(s, "MODULE 07", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Database Maintenance", Inches(0.55), Inches(1.65), Inches(12.2), Inches(1.2),
   fs=46, bold=True, color=WHITE)
tb(s, "Keep NovaMart's indexes healthy, statistics current,\n"
      "and data pages consistent — before problems appear.",
   Inches(0.55), Inches(3.35), Inches(11.0), Inches(0.9), fs=20, color=GREEN_OK)
chips = [
    ("Index Maintenance", BLUE_ACC, Inches(0.55)),
    ("Update Statistics",  ORANGE,  Inches(4.45)),
    ("DBCC CHECKDB",       STEEL,   Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(4.95), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.02), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Day 2  ·  Morning  ·  1.5 hours", Inches(0.55), Inches(6.72), Inches(6.0), Inches(0.35),
   fs=13, color=MID_GRAY)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN_OK)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 07 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (GREEN_OK, "Explain why index fragmentation degrades query performance"),
    (BLUE_ACC, "Decide when to REBUILD versus REORGANIZE an index based on fragmentation level"),
    (ORANGE,   "Update column statistics to ensure the query optimizer has accurate data"),
    (STEEL,    "Run DBCC CHECKDB to detect and report database consistency errors"),
    (GREEN_OK, "Build a weekly maintenance routine for the NovaMart database"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Index Fragmentation ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
tb(s, "Index Fragmentation", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Fragmentation occurs as INSERT, UPDATE, and DELETE operations scatter data across pages.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Two explanation boxes
rect(s, Inches(0.45), Inches(1.55), Inches(5.95), Inches(4.2), NAVY)
tb(s, "WHAT IS FRAGMENTATION?", Inches(0.45), Inches(1.72), Inches(5.95), Inches(0.4),
   fs=17, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
frag_points = [
    "Index pages are stored out of logical order",
    "SQL Server must read more pages to satisfy queries",
    "Reads become random I/O instead of sequential I/O",
    "Worst on tables with frequent inserts and deletes",
    "Measured as avg_fragmentation_in_percent in DMVs",
]
for j, pt in enumerate(frag_points):
    tb(s, f"·  {pt}", Inches(0.65), Inches(2.28) + j * Inches(0.55),
       Inches(5.55), Inches(0.42), fs=14, color=WHITE)

rect(s, Inches(6.78), Inches(1.55), Inches(6.1), Inches(4.2), RGBColor(0x05, 0x12, 0x2A))
tb(s, "HOW TO DETECT IT", Inches(6.78), Inches(1.72), Inches(6.1), Inches(0.4),
   fs=17, bold=True, color=BLUE_ACC, align=PP_ALIGN.CENTER)
detect_sql = (
    "SELECT\n"
    "  OBJECT_NAME(ips.object_id) AS TableName,\n"
    "  i.name AS IndexName,\n"
    "  ips.avg_fragmentation_in_percent,\n"
    "  ips.page_count\n"
    "FROM sys.dm_db_index_physical_stats\n"
    "  (DB_ID(), NULL, NULL, NULL, 'LIMITED') ips\n"
    "JOIN sys.indexes i\n"
    "  ON ips.object_id = i.object_id\n"
    "  AND ips.index_id = i.index_id\n"
    "WHERE ips.avg_fragmentation_in_percent > 5\n"
    "  AND ips.page_count > 100\n"
    "ORDER BY ips.avg_fragmentation_in_percent DESC;"
)
code(s, detect_sql, Inches(6.9), Inches(2.05), Inches(5.9), Inches(3.5), fs=11)
tb(s, "Rule of thumb:  5–30% → REORGANIZE  ·  > 30% → REBUILD  ·  < 5% → skip",
   Inches(0.45), Inches(5.88), Inches(12.45), Inches(0.42), fs=14, bold=True, color=DARK_TEXT)
footer(s, 3)


# ─── Slide 4: Rebuild vs Reorganize ──────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "REBUILD vs REORGANIZE", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Choose based on fragmentation level, available maintenance window, and edition of SQL Server.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# REORGANIZE
rect(s, Inches(0.45), Inches(1.55), Inches(5.85), Inches(3.85), STEEL)
tb(s, "REORGANIZE", Inches(0.45), Inches(1.72), Inches(5.85), Inches(0.4),
   fs=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "5–30% fragmentation", Inches(0.45), Inches(2.18), Inches(5.85), Inches(0.28),
   fs=14, italic=True, color=RGBColor(0xC5, 0xD8, 0xF0), align=PP_ALIGN.CENTER)
reorg_points = [
    "Defragments leaf level pages in-place",
    "Online operation — no table lock",
    "Slower than REBUILD for heavy fragmentation",
    "Statistics NOT automatically updated",
    "Works on individual or all indexes",
]
for j, pt in enumerate(reorg_points):
    tb(s, f"·  {pt}", Inches(0.65), Inches(2.55) + j * Inches(0.52),
       Inches(5.45), Inches(0.38), fs=13, color=WHITE)
reorg_sql = "ALTER INDEX ALL ON dbo.Employee\nREORGANIZE;"
code(s, reorg_sql, Inches(0.45), Inches(5.55), Inches(5.85), Inches(0.78), fs=12)

# REBUILD
rect(s, Inches(6.93), Inches(1.55), Inches(5.95), Inches(3.85), NAVY)
tb(s, "REBUILD", Inches(6.93), Inches(1.72), Inches(5.95), Inches(0.4),
   fs=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "> 30% fragmentation", Inches(6.93), Inches(2.18), Inches(5.95), Inches(0.28),
   fs=14, italic=True, color=MID_GRAY, align=PP_ALIGN.CENTER)
rebuild_points = [
    "Drops and recreates the index from scratch",
    "Offline by default (table locked) — Enterprise: online",
    "Updates statistics automatically (FULLSCAN)",
    "Resets fill factor on all pages",
    "Best fragmentation reduction",
]
for j, pt in enumerate(rebuild_points):
    tb(s, f"·  {pt}", Inches(7.13), Inches(2.55) + j * Inches(0.52),
       Inches(5.55), Inches(0.38), fs=13, color=WHITE)
rebuild_sql = "ALTER INDEX ALL ON dbo.Employee\nREBUILD WITH (ONLINE = OFF);"
code(s, rebuild_sql, Inches(6.93), Inches(5.55), Inches(5.95), Inches(0.78), fs=12)

tb(s, "Rebuild all indexes in NovaMart:  ALTER INDEX ALL ON <table> REBUILD;  — run per table.",
   Inches(0.45), Inches(6.60), Inches(12.45), Inches(0.28), fs=12, italic=True, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: Update Statistics & DBCC CHECKDB ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Update Statistics & DBCC CHECKDB", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Statistics tell the optimizer how data is distributed. CHECKDB validates physical integrity.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

stats_sql = (
    "-- Update statistics for one table\n"
    "UPDATE STATISTICS dbo.Employee\n"
    "WITH FULLSCAN;  -- scan 100% of rows\n\n"
    "-- Update all statistics in a database\n"
    "EXEC sp_updatestats;\n\n"
    "-- Check when statistics were last updated\n"
    "SELECT\n"
    "  s.name, STATS_DATE(s.object_id, s.stats_id) AS LastUpdated\n"
    "FROM sys.stats s\n"
    "WHERE OBJECT_NAME(s.object_id) = 'Employee'\n"
    "ORDER BY LastUpdated DESC;"
)
code(s, stats_sql, Inches(0.45), Inches(1.52), Inches(5.75), Inches(4.35), fs=12)

checkdb_sql = (
    "-- Run full integrity check (takes time on large DBs)\n"
    "DBCC CHECKDB (NovaMart) WITH NO_INFOMSGS;\n\n"
    "-- Fast check — structure only, no row-level\n"
    "DBCC CHECKDB (NovaMart)\n"
    "WITH PHYSICAL_ONLY, NO_INFOMSGS;\n\n"
    "-- Check a single table\n"
    "DBCC CHECKTABLE ('dbo.Employee');\n\n"
    "-- What to look for in output:\n"
    "-- 0 allocation errors, 0 consistency errors\n"
    "-- CHECKDB found 0 allocation errors and\n"
    "-- 0 consistency errors in database 'NovaMart'."
)
code(s, checkdb_sql, Inches(6.72), Inches(1.52), Inches(6.16), Inches(4.35), fs=12)

tb(s, "STATISTICS", Inches(0.45), Inches(1.48), Inches(5.75), Inches(0.3),
   fs=15, bold=True, color=GREEN_OK)
tb(s, "DBCC CHECKDB", Inches(6.72), Inches(1.48), Inches(6.16), Inches(0.3),
   fs=15, bold=True, color=GREEN_OK)

notes = [
    "Run sp_updatestats after every large data load",
    "DBCC CHECKDB weekly minimum — overnight for large databases",
    "If CHECKDB reports errors: do NOT restart the service — escalate immediately",
    "PHYSICAL_ONLY is faster — use for frequent checks; full CHECKDB weekly",
]
for i, note in enumerate(notes):
    col = [GREEN_OK, BLUE_ACC, RED_ERR, STEEL][i]
    rect(s, Inches(0.45), Inches(5.90) + i * Inches(0.28), Inches(0.25), Inches(0.25), col)
    tb(s, note, Inches(0.85), Inches(5.90) + i * Inches(0.28), Inches(12.05), Inches(0.3),
       fs=13, color=DARK_TEXT)
footer(s, 5)


# ─── Slide 6: Maintenance Best Practices ─────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "NovaMart Weekly Maintenance Schedule", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=34, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "A predictable maintenance routine reduces emergency incidents and keeps queries fast.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

schedule = [
    ("Sunday 02:00",  ORANGE,   "Full Backup",
     "BACKUP DATABASE NovaMart ... WITH FORMAT, COMPRESSION"),
    ("Daily 00:00",   BLUE_ACC, "Differential Backup",
     "BACKUP DATABASE NovaMart ... WITH DIFFERENTIAL, COMPRESSION"),
    ("Every 30 min",  GREEN_OK, "Log Backup",
     "BACKUP LOG NovaMart ... — prevents log file growth in FULL recovery"),
    ("Sunday 03:00",  STEEL,    "Index Rebuild",
     "ALTER INDEX ALL ON <each table> REBUILD — after full backup"),
    ("Daily 01:00",   AMBER,    "Update Statistics",
     "EXEC sp_updatestats — after differential backup"),
    ("Sunday 04:00",  NAVY,     "DBCC CHECKDB",
     "DBCC CHECKDB (NovaMart) WITH NO_INFOMSGS — review output"),
]
for i, (time, col, task, detail) in enumerate(schedule):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.65) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(5.95), Inches(1.45), col)
    tb(s, time, x + Inches(0.18), y + Inches(0.1), Inches(2.0), Inches(0.28),
       fs=12, color=WHITE, italic=True)
    tb(s, task, x + Inches(0.18), y + Inches(0.38), Inches(5.58), Inches(0.32),
       fs=16, bold=True, color=WHITE)
    tb(s, detail, x + Inches(0.18), y + Inches(0.78), Inches(5.58), Inches(0.48),
       fs=12, color=WHITE)
footer(s, 6)


# ─── Slide 7: Lab 07 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), GREEN_OK)
tb(s, "LAB 07", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Maintain the NovaMart Database", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~45 min  ·  Tool: SSMS  ·  Database: NovaMart",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Query sys.dm_db_index_physical_stats — identify fragmentation levels in NovaMart",
    "Identify any index with avg_fragmentation_in_percent > 5",
    "Run REORGANIZE on indexes with 5–30% fragmentation",
    "Run REBUILD on indexes with > 30% fragmentation",
    "Run UPDATE STATISTICS dbo.Employee WITH FULLSCAN",
    "Run EXEC sp_updatestats to update all tables",
    "Run DBCC CHECKDB (NovaMart) WITH NO_INFOMSGS",
    "Confirm: '0 allocation errors and 0 consistency errors'",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.47) * i, Inches(0.32), Inches(0.32), GREEN_OK)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.47) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.47) * i, Inches(6.8), Inches(0.36),
       fs=13, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "Fragmentation detected via DMV",
    "REORGANIZE/REBUILD completed",
    "sp_updatestats reports updates",
    "CHECKDB reports 0 errors",
    "All commands run without error",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.70) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)
footer(s, 7)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-07-Database-Maintenance.pptx")
prs.save(out)
print(f"Saved: {out}")
