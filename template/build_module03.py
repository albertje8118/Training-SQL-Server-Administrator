from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "03"
MOD_TITLE = "Database Creation & Management"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def code(slide, sql, l, t, w, h, fs=13):
    """Dark code block for T-SQL snippets."""
    rect(slide, l, t, w, h, CODE_BG)
    box = slide.shapes.add_textbox(l + Inches(0.18), t + Inches(0.12),
                                   w - Inches(0.36), h - Inches(0.24))
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    r = p.add_run()
    r.text = sql
    r.font.size = Pt(fs)
    r.font.color.rgb = RGBColor(0xA8, 0xD8, 0xFF)
    r.font.name = "Consolas"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), GREEN_OK)
tb(s, "MODULE 03", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Database Creation\n& Management", Inches(0.55), Inches(1.65),
   Inches(12.2), Inches(1.6), fs=46, bold=True, color=WHITE)
tb(s, "Learn how SQL Server stores data in files, how to create databases using\nboth SSMS and T-SQL, and how to manage size and growth settings.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=20, color=GREEN_OK)
chips = [
    ("MDF · NDF · LDF",        ORANGE,   Inches(0.55)),
    ("SSMS & T-SQL Creation",  BLUE_ACC, Inches(4.45)),
    ("Filegroups & Properties", STEEL,   Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(5.1), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.17), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN_OK)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 03 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (GREEN_OK, "Explain the role of MDF, NDF, and LDF files in a SQL Server database"),
    (BLUE_ACC, "Create a database using the SSMS GUI with defined file sizes and growth"),
    (ORANGE,   "Create a database using T-SQL CREATE DATABASE with explicit file options"),
    (STEEL,    "Describe how filegroups organize data files and why they are useful"),
    (GREEN_OK, "Modify database properties including auto-growth, compatibility level, and recovery model"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Database File Structure ────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Database File Structure", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Every user database contains at least two files — one for data, one for the transaction log.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
files = [
    (NAVY,     ".MDF",  "Primary Data File",
     ["Contains all schema objects: tables, indexes, views",
      "Every database has exactly one MDF",
      "Stores internal system catalog alongside user data",
      "Location: typically on a fast data drive"]),
    (BLUE_ACC, ".NDF",  "Secondary Data File",
     ["Optional — used to span data across multiple drives",
      "Assigned to filegroups (not PRIMARY by default)",
      "Useful for very large databases or I/O distribution",
      "One or more NDF files per database"]),
    (ORANGE,   ".LDF",  "Transaction Log File",
     ["Records every transaction for rollback and recovery",
      "NEVER on the same drive as .MDF (I/O contention)",
      "Size depends on recovery model and backup frequency",
      "Grows when not truncated by log backup"]),
]
for i, (col, ext, name, bullets) in enumerate(files):
    x = Inches(0.45) + i * Inches(4.22)
    rect(s, x, Inches(1.72), Inches(3.95), Inches(4.75), col)
    tb(s, ext, x, Inches(1.88), Inches(3.95), Inches(0.52),
       fs=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, name, x, Inches(2.48), Inches(3.95), Inches(0.35),
       fs=15, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.2), Inches(2.88), Inches(3.55), Inches(0.02), WHITE)
    for j, b in enumerate(bullets):
        tb(s, f"· {b}", x + Inches(0.25), Inches(3.08) + Inches(0.72) * j,
           Inches(3.55), Inches(0.52), fs=13, color=WHITE)
footer(s, 3)


# ─── Slide 4: Create Database via SSMS ───────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), NAVY)
tb(s, "Creating a Database via SSMS", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Right-click → New Database is the fastest way; set size and growth before clicking OK.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
steps = [
    (NAVY,     "Object Explorer", "Expand the instance → right-click Databases → New Database"),
    (BLUE_ACC, "Database Name",   "Enter the database name (e.g., TrainingDB) — no spaces recommended"),
    (GREEN_OK, "Files Page",      "Review the data file and log file entries — set initial size"),
    (ORANGE,   "Initial Size",    "Data: start at 256 MB+; Log: 128 MB+ — avoid frequent auto-growth"),
    (STEEL,    "Auto-Growth",     "Set data growth to 128 MB (not percent); log to 64 MB"),
    (NAVY,     "File Path",       "Verify the file paths point to your data and log drives"),
    (BLUE_ACC, "Options Page",    "Set Recovery Model (Simple for training), Compatibility Level 150+"),
    (GREEN_OK, "OK",              "Click OK — database appears in Object Explorer immediately"),
]
for i, (col, title, desc) in enumerate(steps):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.72) + row_idx * Inches(1.18)
    rect(s, x, y, Inches(0.5), Inches(0.88), col)
    tb(s, str(i + 1), x, y + Inches(0.22), Inches(0.5), Inches(0.35),
       fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.5), y, Inches(5.6), Inches(0.88), WHITE)
    tb(s, title, x + Inches(0.65), y + Inches(0.08), Inches(5.1), Inches(0.3),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.65), y + Inches(0.42), Inches(5.1), Inches(0.35),
       fs=13, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: Create Database via T-SQL ──────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Creating a Database via T-SQL", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "T-SQL gives you full control over file placement, size, and growth — and is scriptable and repeatable.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

sql_text = (
    "CREATE DATABASE TrainingDB\n"
    "ON PRIMARY\n"
    "(\n"
    "    NAME = TrainingDB_data,\n"
    "    FILENAME = 'D:\\SQLData\\TrainingDB.mdf',\n"
    "    SIZE = 256MB,\n"
    "    MAXSIZE = UNLIMITED,\n"
    "    FILEGROWTH = 128MB\n"
    ")\n"
    "LOG ON\n"
    "(\n"
    "    NAME = TrainingDB_log,\n"
    "    FILENAME = 'E:\\SQLLogs\\TrainingDB.ldf',\n"
    "    SIZE = 128MB,\n"
    "    MAXSIZE = 2048MB,\n"
    "    FILEGROWTH = 64MB\n"
    ");"
)
code(s, sql_text, Inches(0.5), Inches(1.62), Inches(6.55), Inches(5.0), fs=13)

annotations = [
    (NAVY,     "ON PRIMARY",     "Assigns data file to the PRIMARY filegroup"),
    (BLUE_ACC, "NAME",           "Logical name used by SQL Server internally"),
    (GREEN_OK, "FILENAME",       "Physical path on disk — use separate drives"),
    (ORANGE,   "SIZE",           "Initial allocated size — set this generously"),
    (STEEL,    "FILEGROWTH",     "Use MB not % — predictable and auditable"),
    (NAVY,     "LOG ON",         "Separate clause required for the .ldf file"),
]
for i, (col, kw, note) in enumerate(annotations):
    y = Inches(1.62) + i * Inches(0.78)
    rect(s, Inches(7.4), y, Inches(0.38), Inches(0.58), col)
    tb(s, kw, Inches(7.9), y + Inches(0.04), Inches(2.0), Inches(0.28),
       fs=14, bold=True, color=DARK_TEXT)
    tb(s, note, Inches(7.9), y + Inches(0.32), Inches(4.85), Inches(0.28),
       fs=12, color=MID_GRAY)
tb(s, "Tip: run USE TrainingDB; GO after creation to confirm the database is accessible.",
   Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3), fs=13, italic=True, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: Filegroups & Database Properties ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "Filegroups & Database Properties", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Filegroups are logical containers for data files. Key properties control behaviour and recoverability.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Left: Filegroups
rect(s, Inches(0.45), Inches(1.62), Inches(5.95), Inches(4.8), NAVY)
tb(s, "FILEGROUPS", Inches(0.45), Inches(1.78), Inches(5.95), Inches(0.38),
   fs=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
fg_points = [
    "PRIMARY is the default filegroup — contains system objects",
    "Create additional filegroups to separate user tables",
    "Assign data files (.mdf / .ndf) to a filegroup",
    "Set a filegroup as DEFAULT to redirect new objects",
    "Example: separate OLTP tables from archive tables",
    "Useful for partial database restores and availability",
]
for j, pt in enumerate(fg_points):
    tb(s, f"· {pt}", Inches(0.68), Inches(2.32) + Inches(0.55) * j,
       Inches(5.5), Inches(0.42), fs=13, color=WHITE)

# Right: Database Properties
rect(s, Inches(6.75), Inches(1.62), Inches(6.13), Inches(4.8), STEEL)
tb(s, "KEY DATABASE PROPERTIES", Inches(6.75), Inches(1.78), Inches(6.13), Inches(0.38),
   fs=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
props = [
    ("Recovery Model",        "SIMPLE (no log backups) / FULL / BULK_LOGGED"),
    ("Compatibility Level",   "Controls which T-SQL features are available"),
    ("Auto Shrink",           "DISABLE — causes severe fragmentation"),
    ("Auto Close",            "DISABLE — forces re-open on every connection"),
    ("Auto Create Statistics","ON — let SQL Server manage query stats"),
    ("Page Verify",           "CHECKSUM — detects disk-level corruption"),
]
for j, (prop, note) in enumerate(props):
    y = Inches(2.32) + j * Inches(0.62)
    tb(s, prop, Inches(6.95), y, Inches(2.5), Inches(0.28),
       fs=13, bold=True, color=WHITE)
    tb(s, note, Inches(9.55), y, Inches(3.2), Inches(0.28),
       fs=12, color=RGBColor(0xC5, 0xD8, 0xF0))
footer(s, 6)


# ─── Slide 7: Lab 03 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), GREEN_OK)
tb(s, "LAB 03", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Create, Modify & Delete Databases", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=32, bold=True, color=WHITE)
tb(s, "Duration: ~45 min  ·  Tools: SSMS + T-SQL  ·  No external dependencies",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "In SSMS: right-click Databases → New Database → create 'LabDB_GUI'",
    "Set initial data file size to 64 MB, log to 32 MB, growth to 16 MB",
    "Open New Query → run CREATE DATABASE LabDB_TSQL with explicit file paths",
    "Expand both new databases — confirm they appear in Object Explorer",
    "Right-click LabDB_GUI → Properties → Options → change Recovery Model to Simple",
    "Run: ALTER DATABASE LabDB_TSQL MODIFY FILE (NAME=..., SIZE=128MB)",
    "Verify with: SELECT name, size*8/1024 AS SizeMB FROM sys.master_files",
    "Right-click LabDB_GUI → Delete → confirm removal. Verify it disappears.",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.45) * i, Inches(0.32), Inches(0.32), GREEN_OK)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.45) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.45) * i, Inches(6.8), Inches(0.36),
       fs=13, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.08), RGBColor(0x0A, 0x14, 0x32))
tb(s, "VALIDATION", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "Both databases visible in Object Explorer",
    "sys.master_files shows correct sizes",
    "Recovery model updated in Properties",
    "ALTER DATABASE changes file size",
    "Deleted DB no longer appears",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.70) * i,
       Inches(3.88), Inches(0.52), fs=13, color=WHITE)
tb(s, "Clean up: DELETE LabDB_TSQL after confirming all steps.",
   Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3), fs=12, italic=True, color=MID_GRAY)
footer(s, 7)


# ─── Slide 8: Module Summary ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), GREEN_OK)
tb(s, "Module 03 Summary", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "You can now create, configure, and manage SQL Server databases in both SSMS and T-SQL.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
takeaways = [
    (GREEN_OK, "Files are physical, databases are logical",
     "MDF + LDF minimum. NDF for multi-drive expansion. Always separate drives."),
    (BLUE_ACC, "T-SQL gives full control",
     "CREATE DATABASE with explicit SIZE and FILEGROWTH is repeatable and auditable."),
    (ORANGE,   "Set growth in MB, not percent",
     "Percentage-based growth causes unpredictable size spikes under load."),
    (STEEL,    "Recovery Model matters",
     "SIMPLE: no log backup needed. FULL: log backups required to prevent log growth."),
    (NAVY,     "Next: Module 04",
     "Tables & Data Management — schema design, data types, constraints, and CRUD."),
]
for i, (col, title, desc) in enumerate(takeaways):
    col_idx = i % 2
    row_idx = i // 2
    if i == 4:
        x = Inches(3.35)
        y = Inches(5.58)
    else:
        x = Inches(0.45) + col_idx * Inches(6.38)
        y = Inches(1.72) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(5.95), Inches(1.48), col)
    tb(s, title, x + Inches(0.2), y + Inches(0.12), Inches(5.55), Inches(0.35),
       fs=16, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.2), y + Inches(0.55), Inches(5.55), Inches(0.62),
       fs=13, color=WHITE)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-03-Database-Creation-and-Management.pptx")
prs.save(out)
print(f"Saved: {out}")
