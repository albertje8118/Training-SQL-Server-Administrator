from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

NAVY     = RGBColor(0x0D, 0x1B, 0x3E)
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE   = RGBColor(0xE8, 0x50, 0x1A)
BLUE_ACC = RGBColor(0x00, 0x7A, 0xCC)
LIGHT_BG = RGBColor(0xF5, 0xF7, 0xFA)
DARK_TEXT= RGBColor(0x0D, 0x1B, 0x3E)
MID_GRAY = RGBColor(0x8A, 0x93, 0xA6)
GREEN_OK = RGBColor(0x21, 0xA3, 0x66)
STEEL    = RGBColor(0x2E, 0x6E, 0xA6)
CODE_BG  = RGBColor(0x1E, 0x1E, 0x2E)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)
MOD_NUM   = "02"
MOD_TITLE = "Installing & Configuring SQL Server"
MOD_LABEL = f"Module {MOD_NUM}: {MOD_TITLE}"
TOTAL = 8

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]


def bg(slide, color):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


def rect(slide, l, t, w, h, color):
    sh = slide.shapes.add_shape(1, l, t, w, h)
    sh.line.fill.background()
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    return sh


def tb(slide, text, l, t, w, h, fs=18, bold=False, color=WHITE,
       align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(l, t, w, h)
    fr = box.text_frame
    fr.word_wrap = True
    p = fr.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = "Calibri"
    return box


def footer(slide, n):
    tb(slide, MOD_LABEL, Inches(0.4), SLIDE_H - Inches(0.42), Inches(9.5), Inches(0.35),
       fs=11, color=MID_GRAY)
    tb(slide, f"{n} / {TOTAL}", SLIDE_W - Inches(1.6), SLIDE_H - Inches(0.42),
       Inches(1.5), Inches(0.35), fs=11, color=MID_GRAY, align=PP_ALIGN.RIGHT)


# ─── Slide 1: Module Title ────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
rect(s, Inches(0.55), Inches(0.88), Inches(2.2), Inches(0.52), BLUE_ACC)
tb(s, "MODULE 02", Inches(0.55), Inches(0.95), Inches(2.2), Inches(0.38),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Installing &\nConfiguring SQL Server", Inches(0.55), Inches(1.65),
   Inches(12.2), Inches(1.6), fs=44, bold=True, color=WHITE)
tb(s, "From installation planning through TCP/IP connectivity —\nthe steps every administrator needs to get SQL Server production-ready.",
   Inches(0.55), Inches(3.55), Inches(11.0), Inches(0.9), fs=20, color=BLUE_ACC)
chips = [
    ("Pre-Install Planning", ORANGE,   Inches(0.55)),
    ("Installation Steps",   GREEN_OK, Inches(4.45)),
    ("Instance Config",      STEEL,    Inches(8.35)),
]
for label, col, x in chips:
    rect(s, x, Inches(5.1), Inches(3.45), Inches(0.58), col)
    tb(s, label, x, Inches(5.17), Inches(3.45), Inches(0.38),
       fs=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
footer(s, 1)


# ─── Slide 2: Learning Objectives ────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_ACC)
tb(s, "What You Will Learn", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "By the end of Module 02 you will be able to:",
   Inches(0.7), Inches(1.08), Inches(12.0), Inches(0.38), fs=18, italic=True, color=MID_GRAY)
objectives = [
    (BLUE_ACC, "Plan a SQL Server installation — hardware, storage layout, and collation decisions"),
    (ORANGE,   "Run the SQL Server setup wizard and select the right features"),
    (GREEN_OK, "Configure service accounts and set startup types after installation"),
    (STEEL,    "Use SQL Server Configuration Manager to manage services and protocols"),
    (ORANGE,   "Enable TCP/IP and connect to the instance remotely from another machine"),
]
for i, (col, text) in enumerate(objectives):
    y = Inches(1.72) + Inches(0.94) * i
    rect(s, Inches(0.7), y + Inches(0.08), Inches(0.38), Inches(0.38), col)
    tb(s, text, Inches(1.28), y, Inches(11.3), Inches(0.52), fs=18, color=DARK_TEXT)
footer(s, 2)


# ─── Slide 3: Pre-Installation Planning ──────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ORANGE)
tb(s, "Pre-Installation Planning", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "Decisions made before Setup runs are difficult to change later — especially collation.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
planning = [
    (NAVY,     "HARDWARE",
     ["Min: 4 cores, 8 GB RAM, 40 GB disk",
      "Separate drives for data (.mdf) and log (.ldf)",
      "TempDB on its own fast drive for busy systems"]),
    (BLUE_ACC, "COLLATION",
     ["SQL_Latin1_General_CP1_CI_AS is the most common",
      "CI = Case Insensitive  |  AS = Accent Sensitive",
      "Match the application — cannot change after install"]),
    (ORANGE,   "FEATURES",
     ["Database Engine — always required",
      "SQL Server Agent — needed for scheduled jobs",
      "Full-Text / SSRS / SSIS — add only when needed"]),
    (GREEN_OK, "ACCOUNTS",
     ["Use dedicated Windows service accounts (not LocalSystem)",
      "NT SERVICE\\MSSQLSERVER for engine, Agent separate",
      "Principle of least privilege — no local admin needed"]),
]
for i, (col, title, bullets) in enumerate(planning):
    col_num = i % 2
    row_num = i // 2
    x = Inches(0.45) + col_num * Inches(6.35)
    y = Inches(1.72) + row_num * Inches(2.22)
    rect(s, x, y, Inches(5.95), Inches(2.0), col)
    tb(s, title, x, y + Inches(0.15), Inches(5.95), Inches(0.35),
       fs=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        tb(s, f"· {b}", x + Inches(0.18), y + Inches(0.62) + Inches(0.42) * j,
           Inches(5.6), Inches(0.35), fs=13, color=WHITE)
footer(s, 3)


# ─── Slide 4: Installation Step Flow ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), NAVY)
tb(s, "Running the SQL Server Setup Wizard", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "All steps run from the SQL Server installation media — setup.exe.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
steps = [
    (ORANGE,   "1. Installation Center", "Choose 'New SQL Server stand-alone installation'"),
    (BLUE_ACC, "2. Product Key",          "Enter the key or choose Developer (free)"),
    (GREEN_OK, "3. Feature Selection",    "Database Engine Services + SQL Server Agent (minimum)"),
    (STEEL,    "4. Instance Config",      "Default (MSSQLSERVER) or named (SERVER\\INSTANCE)"),
    (ORANGE,   "5. Server Config",        "Set service accounts and collation (SQL_Latin1_General_CP1_CI_AS)"),
    (BLUE_ACC, "6. Auth Mode",            "Windows Auth (recommended) or Mixed Mode"),
    (GREEN_OK, "7. Data Directories",     "Point data root, user DB dir, user log dir, tempDB to correct drives"),
    (NAVY,     "8. Install",              "Progress bar runs — takes 5–15 min. Check summary for errors."),
]
for i, (col, title, desc) in enumerate(steps):
    col_idx = i % 2
    row_idx = i // 2
    x = Inches(0.45) + col_idx * Inches(6.38)
    y = Inches(1.72) + row_idx * Inches(1.18)
    rect(s, x, y, Inches(0.5), Inches(0.88), col)
    tb(s, str(i + 1), x, y + Inches(0.22), Inches(0.5), Inches(0.35),
       fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x + Inches(0.5), y, Inches(5.6), Inches(0.88), WHITE)
    tb(s, title, x + Inches(0.65), y + Inches(0.08), Inches(5.1), Inches(0.3),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, x + Inches(0.65), y + Inches(0.42), Inches(5.1), Inches(0.35),
       fs=13, color=MID_GRAY)
footer(s, 4)


# ─── Slide 5: Instance Configuration ─────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, WHITE)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), GREEN_OK)
tb(s, "Instance Configuration", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "After installation, verify and tune these settings before connecting any application.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)

# Left column
left_items = [
    (NAVY,     "Default vs Named Instance",
     "Default: connect with server name only\nNamed: SERVER\\INSTANCENAME — port may differ"),
    (BLUE_ACC, "Max Server Memory",
     "Limit via SSMS → Server Properties → Memory\nLeave ~10% free for the OS (e.g., 14 GB out of 16 GB)"),
    (GREEN_OK, "Max Degree of Parallelism",
     "MAXDOP controls parallel query threads\nStart at 0 (auto) or 1 per NUMA node rule"),
]
right_items = [
    (ORANGE,   "Cost Threshold for Parallelism",
     "Default 5 is very low — raise to 25–50\nPrevents trivial queries from going parallel"),
    (STEEL,    "Remote Access & Protocols",
     "Shared Memory (local), Named Pipes (legacy),\nTCP/IP (network) — enable in Config Manager"),
    (NAVY,     "Verify with SELECT @@SERVERNAME",
     "Should return the correct host\\instance name\nIf wrong: sp_dropserver / sp_addserver to fix"),
]
for i, (col, title, desc) in enumerate(left_items):
    y = Inches(1.72) + i * Inches(1.55)
    rect(s, Inches(0.45), y, Inches(0.38), Inches(1.25), col)
    tb(s, title, Inches(1.0), y + Inches(0.08), Inches(5.2), Inches(0.35),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, Inches(1.0), y + Inches(0.48), Inches(5.2), Inches(0.55),
       fs=13, color=MID_GRAY)
for i, (col, title, desc) in enumerate(right_items):
    y = Inches(1.72) + i * Inches(1.55)
    rect(s, Inches(6.75), y, Inches(0.38), Inches(1.25), col)
    tb(s, title, Inches(7.3), y + Inches(0.08), Inches(5.2), Inches(0.35),
       fs=15, bold=True, color=DARK_TEXT)
    tb(s, desc, Inches(7.3), y + Inches(0.48), Inches(5.2), Inches(0.55),
       fs=13, color=MID_GRAY)
footer(s, 5)


# ─── Slide 6: SQL Server Configuration Manager ───────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), STEEL)
tb(s, "SQL Server Configuration Manager", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=36, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "The correct tool for service management and protocol configuration — not Windows Services.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
sections = [
    (NAVY,     "SQL Server Services",
     ["Start, stop, pause each SQL Server service",
      "Set startup type: Automatic / Manual / Disabled",
      "View service account — change here, not in Windows Services",
      "Key services: Database Engine, SQL Server Agent, Browser"]),
    (BLUE_ACC, "SQL Server Network Configuration",
     ["Enable/disable protocols per instance",
      "TCP/IP → Properties → IP Addresses: set port 1433",
      "Named Pipes: only needed for legacy apps",
      "Restart the service after any protocol change"]),
    (ORANGE,   "SQL Native Client Configuration",
     ["Client-side protocol order — affects how this machine connects out",
      "Aliases: map a friendly name to an instance\\port",
      "Useful for testing failover or pointing apps to a new server",
      "Rarely changed in standard setups"]),
]
for i, (col, title, bullets) in enumerate(sections):
    x = Inches(0.35) + i * Inches(4.3)
    rect(s, x, Inches(1.72), Inches(3.95), Inches(4.7), col)
    tb(s, title, x, Inches(1.9), Inches(3.95), Inches(0.42),
       fs=17, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, b in enumerate(bullets):
        tb(s, f"· {b}", x + Inches(0.18), Inches(2.5) + Inches(0.72) * j,
           Inches(3.6), Inches(0.52), fs=13, color=WHITE)
footer(s, 6)


# ─── Slide 7: Lab 02 ─────────────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, NAVY)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.08), BLUE_ACC)
rect(s, Inches(0.5), Inches(0.75), Inches(2.0), Inches(0.5), BLUE_ACC)
tb(s, "LAB 02", Inches(0.5), Inches(0.82), Inches(2.0), Inches(0.35),
   fs=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
tb(s, "Configure Services, Protocols & Remote Connectivity", Inches(0.5), Inches(1.42),
   Inches(12.3), Inches(0.62), fs=30, bold=True, color=WHITE)
tb(s, "Duration: ~45 min  ·  Tools: SSMS + Configuration Manager  ·  Assumes: SQL Server Developer installed",
   Inches(0.5), Inches(2.12), Inches(12.3), Inches(0.35), fs=14, color=MID_GRAY)

steps = [
    "Open SQL Server Configuration Manager → verify all services show 'Running'",
    "Set SQL Server Agent startup type to Automatic — restart the service",
    "Under Network Configuration → enable TCP/IP, set port 1433 on IPAll",
    "Restart the Database Engine service to apply protocol changes",
    "In Windows Firewall: add inbound rule for TCP port 1433",
    "From a second machine (or VM): connect in SSMS using the hostname",
    "Verify: SELECT @@SERVERNAME and SELECT @@VERSION both return results",
]
rect(s, Inches(0.5), Inches(2.62), Inches(7.7), Inches(4.05), RGBColor(0x0A, 0x14, 0x32))
tb(s, "STEPS", Inches(0.7), Inches(2.72), Inches(7.0), Inches(0.32),
   fs=14, bold=True, color=BLUE_ACC)
for i, step in enumerate(steps):
    rect(s, Inches(0.7), Inches(3.12) + Inches(0.52) * i, Inches(0.32), Inches(0.32), BLUE_ACC)
    tb(s, str(i + 1), Inches(0.7), Inches(3.14) + Inches(0.52) * i, Inches(0.32), Inches(0.28),
       fs=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    tb(s, step, Inches(1.15), Inches(3.10) + Inches(0.52) * i, Inches(6.8), Inches(0.38),
       fs=13, color=WHITE)

rect(s, Inches(8.55), Inches(2.62), Inches(4.28), Inches(4.05), RGBColor(0x0A, 0x14, 0x32))
tb(s, "EXPECTED OUTCOME", Inches(8.75), Inches(2.72), Inches(3.88), Inches(0.32),
   fs=14, bold=True, color=GREEN_OK)
outcomes = [
    "SQL Server Agent is running on startup",
    "TCP/IP enabled, port 1433 confirmed",
    "Remote SSMS connection succeeds",
    "Firewall rule allows port 1433 traffic",
]
for i, o in enumerate(outcomes):
    tb(s, f"✓  {o}", Inches(8.75), Inches(3.12) + Inches(0.75) * i,
       Inches(3.88), Inches(0.55), fs=13, color=WHITE)
tb(s, "Clean up: no user databases created in this lab.",
   Inches(0.5), Inches(6.75), Inches(12.3), Inches(0.3), fs=12, italic=True, color=MID_GRAY)
footer(s, 7)


# ─── Slide 8: Module Summary ─────────────────────────────────────────────────
s = prs.slides.add_slide(blank)
bg(s, LIGHT_BG)
rect(s, Inches(0), Inches(0), SLIDE_W, Inches(0.06), BLUE_ACC)
tb(s, "Module 02 Summary", Inches(0), Inches(0.38), SLIDE_W, Inches(0.62),
   fs=38, bold=True, color=DARK_TEXT, align=PP_ALIGN.CENTER)
tb(s, "SQL Server is installed, configured, and reachable. You are ready to create databases.",
   Inches(0), Inches(1.05), SLIDE_W, Inches(0.38), fs=17, italic=True, color=MID_GRAY,
   align=PP_ALIGN.CENTER)
takeaways = [
    (BLUE_ACC, "Plan before you install",
     "Collation cannot change after install. Separate drives improve performance."),
    (ORANGE,   "Setup wizard is sequential",
     "Feature selection → instance → service accounts → auth mode → directories."),
    (GREEN_OK, "Configuration Manager is key",
     "Manage services and protocols here — not Windows Services or the registry."),
    (STEEL,    "TCP/IP + firewall = remote access",
     "Enable protocol, set port 1433, open firewall. Then restart the service."),
    (NAVY,     "Next: Module 03",
     "Database Creation & Management — MDF, LDF, filegroups, and T-SQL CREATE DATABASE."),
]
for i, (col, title, desc) in enumerate(takeaways):
    col_idx = i % 2
    row_idx = i // 2
    if i == 4:
        x = Inches(3.35)
        y = Inches(5.58)
    else:
        x = Inches(0.45) + col_idx * Inches(6.38)
        y = Inches(1.72) + row_idx * Inches(1.72)
    rect(s, x, y, Inches(5.95), Inches(1.48), col)
    tb(s, title, x + Inches(0.2), y + Inches(0.12), Inches(5.55), Inches(0.35),
       fs=16, bold=True, color=WHITE)
    tb(s, desc, x + Inches(0.2), y + Inches(0.55), Inches(5.55), Inches(0.62),
       fs=13, color=WHITE)
footer(s, 8)


# ─── Save ─────────────────────────────────────────────────────────────────────
repo_root  = os.path.dirname(os.path.dirname(__file__))
output_dir = os.path.join(repo_root, "contents")
os.makedirs(output_dir, exist_ok=True)
out = os.path.join(output_dir, "Module-02-Installing-and-Configuring-SQL-Server.pptx")
prs.save(out)
print(f"Saved: {out}")
